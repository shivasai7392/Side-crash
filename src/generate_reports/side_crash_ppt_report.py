# PYTHON script
"""
    _summary_

_extended_summary_

Returns:
    _type_: _description_
"""

import os
import time

from src.generate_reports.report_slides.title_slide import TitleSlide
from src.generate_reports.report_slides.cae_quality_slide import CAEQualitySlide
from src.generate_reports.report_slides.executive_slide import ExecutiveSlide
from src.generate_reports.report_slides.cbu_and_barrier_position_slide import CBUAndBarrierPositionSlide
from src.generate_reports.report_slides.biw_kinematics_slide import BIWKinematicsSlide
from src.generate_reports.report_slides.biw_cbu_deformation_slide import BIWCBUDeformationSlide
from src.generate_reports.report_slides.bom_f21_upb_slide import BOMF21UPBSlide
from src.generate_reports.report_slides.biw_stiff_ring_deformation_slide import BIWStiffRingDeformationSlide
from src.generate_reports.report_slides.biw_bplr_def_and_intr_slide import BIWBplrDeformationAndIntrusion
class SideCrashPPTReport():

    def __init__(self,
                windows,
                general_input,
                metadb_2d_input,
                metadb_3d_input,
                template_file,
                twod_images_report_folder,
                threed_images_report_folder,
                ppt_report_folder) -> None:
        self.windows = windows
        self.general_input = general_input
        self.metadb_2d_input = metadb_2d_input
        self.metadb_3d_input = metadb_3d_input
        self.template_file = template_file
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.ppt_report_folder = ppt_report_folder
        self.report_composer = None


    def generate_ppt(self):
        """
        thesis_report_generation [summary]

        [extended_summary]

        Returns:
            [type]: [description]
        """

        self.report_composer = PPTXReportComposer(report_name="Run1",template_pptx=self.template_file)
        self.report_composer.create_prs_obj()
        title_slide = TitleSlide(self.report_composer.prs_obj.slides[0],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        title_slide.edit()
        cae_quality_slide = CAEQualitySlide(self.report_composer.prs_obj.slides[1],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        cae_quality_slide.edit()
        executive_slide = ExecutiveSlide(self.report_composer.prs_obj.slides[2],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        executive_slide.edit()
        cbu_and_barrier_position_slide = CBUAndBarrierPositionSlide(self.report_composer.prs_obj.slides[3],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        cbu_and_barrier_position_slide.edit()
        biw_kinematics_slide = BIWKinematicsSlide(self.report_composer.prs_obj.slides[6],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        biw_kinematics_slide.edit()
        biw_cbu_deformation_slide = BIWCBUDeformationSlide(self.report_composer.prs_obj.slides[7],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        biw_cbu_deformation_slide.edit()
        bom_f21_upb_slide = BOMF21UPBSlide(self.report_composer.prs_obj.slides[8],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        bom_f21_upb_slide.edit()
        biw_stiff_ring_deformation_slide = BIWStiffRingDeformationSlide(self.report_composer.prs_obj.slides[9],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        biw_stiff_ring_deformation_slide.edit()
        biw_bplr_def_and_intr_slide = BIWBplrDeformationAndIntrusion(self.report_composer.prs_obj.slides[10],
                                self.windows,
                                self.general_input,
                                self.metadb_2d_input,
                                self.metadb_3d_input,
                                self.template_file,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder,
                                self.ppt_report_folder)
        biw_bplr_def_and_intr_slide.edit()

        if not os.path.exists(self.ppt_report_folder):
            os.makedirs(self.ppt_report_folder)
        file_name = os.path.join(self.ppt_report_folder,"output.pptx")
        self.report_composer.save_pptx(file_name)

        return 0

class PPTXReportComposer():
    """
    __init__ _summary_

    _extended_summary_

    Args:
        report_name (_type_): _description_
        template_pptx (_type_): _description_
    """
    def __init__(self, report_name, template_pptx):
        self.report_name = report_name
        self.template_pptx = template_pptx
        self.prs_obj = None

    def create_prs_obj(self):
        """ Creates the PPTx report using
        the python-pptx module
        """
        from pptx import Presentation

        # Instantiate
        if not self.prs_obj:
            self.prs_obj = Presentation(self.template_pptx)

        return 0

    def save_pptx(self, pptx_filepath, datestamp=""):
        """ Saves the PPTx at the given filepath
        with the given datestamp

        Args:
            pptx_filepath (str): Absolute path to the pptx file for saving.
            datestamp (str, optional): Date stamp at the bottom right of slide. Defaults to "".

        Returns:
            int: 0 always.
        """
        from pptx.enum.text import PP_ALIGN

        # Get current date if not provided
        if not datestamp:
            datestamp = time.strftime("%B %d, %Y")

        # Set the date in the ppt master slide
        master_slide = self.prs_obj.slide_master
        for shape in master_slide.shapes:
            try:
                if shape.text == "Date_Stamp":
                    shape.text = datestamp
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            except AttributeError:
                continue

        # Save the pptx
        self.prs_obj.save(pptx_filepath)

        return 0
