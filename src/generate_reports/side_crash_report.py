import os
from platform import node
import time


from meta import utils,parts,constants
from meta import windows
from meta import plot2d

class SideCrashReport():

    def __init__(self,windows,general_input,metadb_2d_input,metadb_3d_input,config_folder) -> None:
        self.windows = windows
        self.general_input = general_input
        self.metadb_2d_input = metadb_2d_input
        self.metadb_3d_input = metadb_3d_input
        self.config_folder = config_folder
        self.template_file = os.path.join(self.config_folder,"res",self.general_input.source_template_file_directory.replace("/","",1),self.general_input.source_template_file_name).replace("\\",os.sep)
        self.intrusion_areas = {"ROW 1":{},"ROW 2":{}}
        self.get_reporting_folders()

    def get_reporting_folders(self):
        self.twod_images_report_folder = os.path.join(self.config_folder,"res",os.path.dirname(self.general_input.report_directory).replace("/","",1),"2d-data-images").replace("\\",os.sep)
        self.threed_images_report_folder = os.path.join(self.config_folder,"res",os.path.dirname(self.general_input.report_directory).replace("/","",1),"3d-data-images").replace("\\",os.sep)
        self.threed_videos_report_folder = os.path.join(self.config_folder,"res",os.path.dirname(self.general_input.report_directory).replace("/","",1),"3d-data-videos").replace("\\",os.sep)
        self.excel_bom_report_folder = os.path.join(self.config_folder,"res",os.path.dirname(self.general_input.report_directory).replace("/","",1),"excel-bom").replace("\\",os.sep)
        self.ppt_report_folder = os.path.join(self.config_folder,"res",os.path.dirname(self.general_input.report_directory).replace("/","",1),"reports").replace("\\",os.sep)

        return 0

    def run_process(self):

        self.twod_data_reporting()
        self.threed_data_reporting()
        self.thesis_report_generation()
        return 0


    def thesis_report_generation(self):
        """
        thesis_report_generation [summary]

        [extended_summary]

        Returns:
            [type]: [description]
        """

        self.report_composer = PPTXReportComposer(report_name="Run1",template_pptx=self.template_file)
        self.report_composer.create_prs_obj()

        utils.MetaCommand('options title off')
        utils.MetaCommand('options axis off')
        self.edit_title_slide(self.report_composer.prs_obj.slides[0])
        self.edit_cae_quality_slide(self.report_composer.prs_obj.slides[1])
        self.edit_executive_slide(self.report_composer.prs_obj.slides[2])
        self.edit_cbu_and_barrier_position_slide(self.report_composer.prs_obj.slides[3])
        self.edit_body_in_white_kinematics_slide(self.report_composer.prs_obj.slides[6])
        self.edit_body_in_white_cbu_deformation_slide(self.report_composer.prs_obj.slides[7])
        self.edit_bill_of_materials_f21_upb(self.report_composer.prs_obj.slides[8])
        self.biw_stiff_ring_deformation(self.report_composer.prs_obj.slides[9])
        if not os.path.exists(self.ppt_report_folder):
            os.makedirs(self.ppt_report_folder)
        file_name = os.path.join(self.ppt_report_folder,"output.pptx")
        self.report_composer.save_pptx(file_name)

        return 0
    def edit_body_in_white_cbu_deformation_slide(self, slide):
        utils.MetaCommand('window maximize "MetaPost"')
        for shape in slide.shapes:
            if shape.name == "Image 1":
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('grstyle scalarfringe disable')
                data = self.metadb_3d_input.critical_sections
                entities = list()
                for _key,value in data.items():
                    if 'hes' in value.keys():
                        prop_names = value['hes']
                        re_props = prop_names.split(",")
                        for re_prop in re_props:
                            entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('view default isometric')
                utils.MetaCommand('options fringebar off')

                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"cbu_deformation".lower()+".jpeg")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('grstyle scalarfringe enable')

            elif shape.name == "Image 2":
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('grstyle scalarfringe enable')
                data = self.metadb_3d_input.critical_sections
                entities = list()
                for _key,value in data.items():
                    if 'hes' in value.keys():
                        prop_names = value['hes']
                        re_props = prop_names.split(",")
                        for re_prop in re_props:
                            entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('view default isometric')
                utils.MetaCommand('options fringebar off')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"cbu_deformation".lower()+".jpeg")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('0:options state variable "serial=0"')
            elif shape.name == "Image 3":
                utils.MetaCommand('add all')
                utils.MetaCommand('add invert')
                utils.MetaCommand('options fringebar on')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"Fringe_bar".lower()+".jpeg")
                utils.MetaCommand('write scalarfringebar png {} '.format(image_path))
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('options fringebar off')

        return 0
    @staticmethod
    def closest(list_of_values, value):
        """
        closest _summary_

        _extended_summary_

        Args:
            lst (_type_): _description_
            K (_type_): _description_

        Returns:
            _type_: _description_
        """
        if type(value) == int:
            nearest_value = list_of_values[min(range(len(list_of_values)), key = lambda i: abs(list_of_values[i]-value))]
            return nearest_value

        else:
            return None

    def biw_stiff_ring_deformation(self,slide):
        """
        biw_stiff_ring_deformation _summary_

        _extended_summary_

        Args:
            slide (_type_): _description_

        Returns:
            _type_: _description_
        """
        from PIL import ImageGrab
        utils.MetaCommand('0:options state variable "serial=1"')
        for shape in slide.shapes:
            if shape.name == "Image 6":
                window_name = self.general_input.biw_stiff_ring_deformation_name
                win = windows.Window(str(window_name), page_id=0)
                layout = win.get_plot_layout()
                utils.MetaCommand('window maximize "{}"'.format(window_name))
                plot_id = 0
                page_id=0
                final_time_variable =  dict(utils.MetaGetVariablesByName("survival-space_final_time"))
                final_time_roof = final_time_variable["survival-space_final_time"]
                final_time_roof_splitting = final_time_roof.split(".")[0]
                plot = plot2d.Plot(plot_id, window_name, page_id)
                curvelist_final_time = plot.get_curves('byname', name ="ROOF_LINE_"+str(final_time_roof_splitting)+"MS")[0]
                final_time_id = curvelist_final_time.id
                curvelist_initial_time = plot.get_curves('byname', name ="ROOF_LINE_"+str(0)+"MS")[0]
                peak_time_variable = dict(utils.MetaGetVariablesByName("peak_time_display"))
                peak_time_value = peak_time_variable["peak_time_display"]
                peak_time = peak_time_value.split(".")[0]
                plot = plot2d.Plot(plot_id, window_name, page_id)
                curves = plot.get_curves('all')
                final_roof_line_list = list()
                for each_curve in curves:
                    ms = each_curve.name.split("_")[2]
                    if 'MS' in ms:
                        ms_replacing = ms.replace('MS',"")
                        final_roof_line_list.append(int(ms_replacing))
                peak_time_value = self.closest(final_roof_line_list, int(peak_time))
                peak_time_curve = plot.get_curves('byname', name ="ROOF_LINE_"+str(peak_time_value)+"MS")
                for each_peak_time_curve in peak_time_curve:
                    peak_time_id = each_peak_time_curve.id
                initial_time_id = curvelist_initial_time.id
                title = plot2d.Title(plot_id, window_name, page_id)
                plot = title.get_plot()
                plot.activate()
                utils.MetaCommand('xyplot plotactive "{}" 0'.format(window_name))
                utils.MetaCommand('xyplot rlayout "{}" 1'.format(window_name))
                utils.MetaCommand('xyplot curve visible and "{}" {} {},{}'.format(window_name,initial_time_id,peak_time_id, final_time_id))
                utils.MetaCommand('xyplot curve set style "{}" {} 9'.format(window_name, initial_time_id))
                utils.MetaCommand('xyplot curve set style "{}" {} 5'.format(window_name,peak_time_id))
                utils.MetaCommand('xyplot curve select "{}" all'.format(window_name))
                utils.MetaCommand('xyplot axisoptions yaxis active "{}" 0 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions ylabel font "{}" 0 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions labels yfont "{}" 0 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" 0 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions xaxis active "{}" 0 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions xlabel font "{}" 0 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions labels xfont "{}" 0 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot plotoptions title font "{}" 0 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(window_name, plot.id))

                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+title.get_text().lower()+".png").replace("&","and")
                if not os.path.exists(os.path.dirname(image_path)):
                    print(os.path.dirname(image_path))
                    os.makedirs(os.path.dirname(image_path))
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('xyplot rlayout "{}" {}'.format(window_name, layout))
            elif shape.name == "Image 7":
                window_name = self.general_input.biw_stiff_ring_deformation_name
                win = windows.Window(str(window_name), page_id=0)
                layout = win.get_plot_layout()
                utils.MetaCommand('window maximize "{}"'.format(window_name))
                plot_id = 1
                page_id=0
                final_time_variable =  dict(utils.MetaGetVariablesByName("survival-space_final_time"))
                final_time_roof = final_time_variable["survival-space_final_time"]
                final_time_roof_splitting = final_time_roof.split(".")[0]
                plot = plot2d.Plot(plot_id, window_name, page_id)
                curvelist_final_time = plot.get_curves('byname', name ="SIDE_SILL_"+str(final_time_roof_splitting)+"MS")
                for each_curvelist_final_time in curvelist_final_time:
                    final_time_id = each_curvelist_final_time.id
                curvelist_initial_time = plot.get_curves('byname', name ="SIDE_SILL_"+str(0)+"MS")

                peak_time_variable = dict(utils.MetaGetVariablesByName("peak_time_display"))
                peak_time_value = peak_time_variable["peak_time_display"]
                peak_time = peak_time_value.split(".")[0]
                plot = plot2d.Plot(plot_id, window_name, page_id)
                curves = plot.get_curves('all')
                final_roof_line_list = list()
                for each_curve in curves:
                    ms = each_curve.name.split("_")[2]
                    if 'MS' in ms:
                        ms_replacing = ms.replace('MS',"")
                        final_roof_line_list.append(int(ms_replacing))
                peak_time_value = self.closest(final_roof_line_list, int(peak_time))
                peak_time_curve = plot.get_curves('byname', name ="SIDE_SILL_"+str(peak_time_value)+"MS")
                for each_peak_time_curve in peak_time_curve:
                    peak_time_id = each_peak_time_curve.id

                for each_curvelist_initial_time in curvelist_initial_time:
                    initial_time_id = each_curvelist_initial_time.id

                title = plot2d.Title(plot_id, window_name, page_id)
                plot = title.get_plot()
                plot.activate()
                utils.MetaCommand('xyplot plotactive "{}" 1'.format(window_name))
                utils.MetaCommand('xyplot rlayout "{}" 1'.format(window_name))
                utils.MetaCommand('xyplot curve visible and "{}" {} {},{}'.format(window_name,initial_time_id,peak_time_id, final_time_id))
                utils.MetaCommand('xyplot curve set style "{}" {} 9'.format(window_name, initial_time_id))
                utils.MetaCommand('xyplot curve set style "{}" {} 5'.format(window_name,peak_time_id))
                utils.MetaCommand('xyplot curve select "{}" all'.format(window_name))
                utils.MetaCommand('xyplot axisoptions yaxis active "{}" 1 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions axyrange "{}" 1 0 -805 -770'.format(window_name))
                utils.MetaCommand('xyplot axisoptions ylabel font "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions labels yfont "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" 1 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions xaxis active "{}" 1 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions xlabel font "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions labels xfont "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot plotoptions title font "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(window_name, plot.id))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+title.get_text().lower()+".png").replace("&","and")
                if not os.path.exists(os.path.dirname(image_path)):
                    print(os.path.dirname(image_path))
                    os.makedirs(os.path.dirname(image_path))
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('xyplot rlayout "{}" {}'.format(window_name, layout))
            elif shape.name == "Image 4":
                data = self.metadb_3d_input.critical_sections["f21_upb_inner"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('options fringebar off')
                utils.MetaCommand('view default right')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"deformed_peak_intrusion_inner".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Image 5":
                data = self.metadb_3d_input.critical_sections["f21_upb_inner"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('options fringebar off')
                utils.MetaCommand('grstyle deform off')
                utils.MetaCommand('view default right')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"undeformed_peak_intrusion_inner".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('grstyle deform on')
            elif shape.name == "Image 2":
                data = self.metadb_3d_input.critical_sections["f21_upb_outer"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('options fringebar off')
                utils.MetaCommand('grstyle deform on')
                utils.MetaCommand('view default left')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"deformed_peak_intrusion_outer".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path, view = "left")
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Image 3":
                data = self.metadb_3d_input.critical_sections["f21_upb_outer"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('options fringebar off')
                utils.MetaCommand('grstyle deform off')
                utils.MetaCommand('view default left')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"undeformed_peak_intrusion_outer".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path, view = "left")
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Image 1":
                utils.MetaCommand('add all')
                utils.MetaCommand('add invert')
                utils.MetaCommand('options fringebar on')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"Fringe_bar".lower()+".jpeg")
                utils.MetaCommand('write scalarfringebar png {} '.format(image_path))
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('options fringebar off')

        return 0

    def edit_bill_of_materials_f21_upb(self, slide):

        from PIL import ImageGrab
        from pptx.util import Pt

        utils.MetaCommand('0:options state original')
        for shape in slide.shapes:
            if shape.name == "Image 2":
                data = self.metadb_3d_input.critical_sections["f21_upb_inner"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('view default right')
                utils.MetaCommand('options fringebar off')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"f21_upb_inner".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            if shape.name == "Image 1":
                data = self.metadb_3d_input.critical_sections["f21_upb_outer"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('view default left')
                utils.MetaCommand('options fringebar off')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"f21_upb_inner".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Table 1":
                data = self.metadb_3d_input.critical_sections["f21_upb_outer"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                table_obj = shape.table
                for id,prop in enumerate(entities[:15]):
                    part = parts.Part(id=prop.id, type=constants.PSHELL, model_id=0)
                    materials = part.get_materials('all')

                    self.report_composer.add_row(table_obj)
                    prop_row = table_obj.rows[id+1]

                    text_frame = prop_row.cells[0].text_frame
                    font = text_frame.paragraphs[0].font
                    font.size = Pt(8)
                    text_frame.paragraphs[0].text = str(prop.id)

                    text_frame_name = prop_row.cells[1].text_frame
                    font_name = text_frame_name.paragraphs[0].font
                    font_name.size = Pt(8)
                    text_frame_name.paragraphs[0].text = str(prop.name)

                    text_frame_material = prop_row.cells[2].text_frame
                    font_material = text_frame_material.paragraphs[0].font
                    font_material.size = Pt(8)
                    for each_material in materials:
                        text_frame_material.paragraphs[0].text = str(each_material.name)

                    text_frame_thickness = prop_row.cells[3].text_frame
                    font_thickness = text_frame_thickness.paragraphs[0].font
                    font_thickness.size = Pt(8)
                    thickness = round(prop.shell_thick,1)
                    text_frame_thickness.paragraphs[0].text = str(thickness)


            elif shape.name == "Table 2":
                data = self.metadb_3d_input.critical_sections["f21_upb_outer"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))

                entities_all = []
                for each_entity in entities[15:]:
                    if str(each_entity.id).startswith("2"):
                        entities_all.append(each_entity)

                table_obj = shape.table
                for id,prop in enumerate(entities_all):

                    part = parts.Part(id=prop.id, type=constants.PSHELL, model_id=0)
                    materials = part.get_materials('all')

                    self.report_composer.add_row(table_obj)
                    prop_row = table_obj.rows[id+1]
                    text_frame = prop_row.cells[0].text_frame
                    font = text_frame.paragraphs[0].font
                    font.size = Pt(8)
                    text_frame.paragraphs[0].text = str(prop.id)

                    text_frame_name = prop_row.cells[1].text_frame
                    font_name = text_frame_name.paragraphs[0].font
                    font_name.size = Pt(8)
                    text_frame_name.paragraphs[0].text = str(prop.name)

                    text_frame_material = prop_row.cells[2].text_frame
                    font_material = text_frame_material.paragraphs[0].font
                    font_material.size = Pt(8)
                    for each_material in materials:
                        text_frame_material.paragraphs[0].text = str(each_material.name)

                    text_frame_thickness = prop_row.cells[3].text_frame
                    font_thickness = text_frame_thickness.paragraphs[0].font
                    font_thickness.size = Pt(8)
                    thickness = round(prop.shell_thick,1)
                    text_frame_thickness.paragraphs[0].text = str(thickness)
            elif shape.name == "Table 3":
                data = self.metadb_3d_input.critical_sections["f21_upb_inner"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))

                entities_all = []
                for each_entity in entities:
                    if str(each_entity.id).startswith("2"):
                        entities_all.append(each_entity)

                table_obj = shape.table
                for id,prop in enumerate(entities_all):

                    part = parts.Part(id=prop.id, type=constants.PSHELL, model_id=0)
                    materials = part.get_materials('all')

                    self.report_composer.add_row(table_obj)
                    prop_row = table_obj.rows[id+1]

                    text_frame = prop_row.cells[0].text_frame
                    font = text_frame.paragraphs[0].font
                    font.size = Pt(8)
                    text_frame.paragraphs[0].text = str(prop.id)

                    text_frame_name = prop_row.cells[1].text_frame
                    font_name = text_frame_name.paragraphs[0].font
                    font_name.size = Pt(8)
                    text_frame_name.paragraphs[0].text = str(prop.name)

                    text_frame_material = prop_row.cells[2].text_frame
                    font_material = text_frame_material.paragraphs[0].font
                    font_material.size = Pt(8)
                    for each_material in materials:
                        text_frame_material.paragraphs[0].text = str(each_material.name)

                    text_frame_thickness = prop_row.cells[3].text_frame
                    font_thickness = text_frame_thickness.paragraphs[0].font
                    font_thickness.size = Pt(8)
                    thickness = round(prop.shell_thick,1)
                    text_frame_thickness.paragraphs[0].text = str(thickness)

            elif shape.name == "TextBox 1":
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = "OUTER"
            elif shape.name == "TextBox 2":
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = "INNER"
        return 0

    def edit_body_in_white_kinematics_slide(self, slide):
        from PIL import ImageGrab,Image
        utils.MetaCommand('window maximize "MetaPost"')

        for shape in slide.shapes:
            if shape.name == "Image 1":
                utils.MetaCommand('add all')
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('view default front')
                utils.MetaCommand('color fringebar scalarset Critical')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"model_front".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path,rotate = Image.ROTATE_90)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('color fringebar scalarset default')
            elif shape.name == "Image 2":
                utils.MetaCommand('add all')
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('view default top')
                utils.MetaCommand('color fringebar scalarset Critical')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"model_top".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path,rotate = Image.ROTATE_90)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('color fringebar scalarset default')
            elif shape.name == "Image 3":
                utils.MetaCommand('add all')
                utils.MetaCommand('view default isometric')
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('view default front')
                utils.MetaCommand('color fringebar scalarset Critical')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"model_front".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path,rotate = Image.ROTATE_90)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('color fringebar scalarset default')
            elif shape.name == "Image 4":
                utils.MetaCommand('add all')
                utils.MetaCommand('view default isometric')
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('view default top')
                utils.MetaCommand('color fringebar scalarset Critical')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"model_top".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path,rotate = Image.ROTATE_90)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('color fringebar scalarset default')

    def edit_cbu_and_barrier_position_slide(self, slide):

        from PIL import ImageGrab,Image
        utils.MetaCommand('window maximize "MetaPost"')

        for shape in slide.shapes:
            if shape.name == "Image 4":
                #data = self.metadb_3d_input.critical_sections["barrier_and_cbu"]
                data = self.metadb_3d_input.critical_sections["cbu"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('0:options state original')
                    utils.MetaCommand('add all')
                    utils.MetaCommand('view default isometric')
                    utils.MetaCommand('view default top')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"cbu".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path,rotate = Image.ROTATE_90)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            if shape.name == "Image 3":
                #data = self.metadb_3d_input.critical_sections["barrier_and_cbu"]
                data = self.metadb_3d_input.critical_sections["cbu"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('0:options state variable "serial=1"')
                    utils.MetaCommand('add all')
                    utils.MetaCommand('view default isometric')
                    utils.MetaCommand('grstyle deform off')
                    utils.MetaCommand('color fringebar scalarset Critical')
                    utils.MetaCommand('view default left')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('grstyle deform on')
                utils.MetaCommand('color fringebar scalarset default')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"cbu_critical".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path,rotate = Image.ROTATE_90)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0

        return 0

    @staticmethod
    def capture_image(window_name,width,height,file_path,plot_id = None,smoothing = False,rotate = None, view = None):
        """
        capture_image _summary_

        _extended_summary_

        Args:
            window_name (_type_): _description_
            width (_type_): _description_
            height (_type_): _description_
            file_path (_type_): _description_
            plot_id (_type_, optional): _description_. Defaults to None.

        Returns:
            _type_: _description_
        """
        from PIL import ImageGrab

        win_obj = windows.Window(window_name, page_id = 0)
        win_obj.set_size((round(width/9525),round(height/9525)))
        if view is not None:
            utils.MetaCommand('view default {}'.format(view))

        if smoothing:
            utils.MetaCommand('write options outputsize sizesmoothscale {},{}'.format(round(width/9525),round(height/9525)))
        else:
            utils.MetaCommand('write options outputsize smoothscale 1')
        if plot_id is not None:
            utils.MetaCommand('clipboard copy plot image "{}" {}'.format(window_name, plot_id))
        else:
            utils.MetaCommand('clipboard copy image "{}"'.format(window_name))

        img = ImageGrab.grabclipboard()
        rgba_img = img.convert("RGBA")
        rgba_data = rgba_img.getdata()
        new_rgba_data = []
        for item in rgba_data:
            if item[0] == 255 and item[1] == 255 and item[2] == 255:
                new_rgba_data.append((255, 255, 255, 0))
            else:
                new_rgba_data.append(item)
        rgba_img.putdata(new_rgba_data)
        if rotate:
            rgba_img.transpose(rotate)
        if not os.path.exists(os.path.dirname(file_path)):
            os.makedirs(os.path.dirname(file_path))
        rgba_img.save(file_path, 'PNG')

        utils.MetaCommand('window maximize {}'.format(window_name))

        return 0

    def edit_executive_slide(self,slide):
        """
        edit_executive_slide _summary_

        _extended_summary_

        Args:
            slide (_type_): _description_

        Returns:
            _type_: _description_
        """
        from PIL import ImageGrab
        from pptx.util import Pt


        window_name = self.general_input.survival_space_window_name
        utils.MetaCommand('window maximize "{}"'.format(window_name))

        shapes = [shape for shape in slide.shapes]
        shapes.sort(key = lambda x:x.name)

        for shape in shapes:
            if shape.name == "Image":
                final_time_curve_name = "SS_{}MS".format(round(float(self.general_input.survival_space_final_time)))
                initial_time_curve_name = "SS_0MS"
                plot_id = 0
                page_id=0
                plot = plot2d.Plot(plot_id, window_name, page_id)
                title = plot.get_title()
                org_name = title.get_text()
                title.set_text("{} 0MS AND 150MS".format(org_name))
                plot.activate()
                final_curve = plot2d.CurvesByName(window_name, final_time_curve_name, 0)[0]
                final_curve.show()
                initial_curve = plot2d.CurvesByName(window_name, initial_time_curve_name, 0)[0]
                initial_curve.show()
                utils.MetaCommand('xyplot axisoptions xaxis active "Survival Space" 0 0')
                utils.MetaCommand('xyplot gridoptions xspace "Survival Space" 0 20')
                utils.MetaCommand('xyplot axisoptions axxrange "Survival Space" 0 0 175 400')
                utils.MetaCommand('xyplot gridoptions line major style "Survival Space" 0 0')
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+title.get_text().lower()+".png")
                self.capture_image(window_name,shape.width,shape.height,image_path,plot_id=plot.id)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                title.set_text(org_name)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Table 1":
                final_time = round(float(self.general_input.survival_space_final_value))
                rows = shape.table.rows
                text_frame_1 = rows[2].cells[0].text_frame
                font = text_frame_1.paragraphs[0].font
                font.size = Pt(11)
                text_frame_2 = rows[2].cells[1].text_frame
                font = text_frame_2.paragraphs[0].font
                font.size = Pt(11)
                text_frame_3 = rows[3].cells[1].text_frame
                font = text_frame_3.paragraphs[0].font
                font.size = Pt(11)
                text_frame_1.paragraphs[0].text = " {}".format(final_time)
                text_frame_2.paragraphs[0].text = str(round(float(self.general_input.survival_space_peak_value)))
                text_frame_3.paragraphs[0].text = str(round(float(self.general_input.peak_time_display_value)))
            elif shape.name == "Image 3":
                data = self.metadb_3d_input.critical_sections["f28_front_door"]
                prop_names = data["hes"]
                erase_box = data["erase_box"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    utils.MetaCommand('add all')
                    utils.MetaCommand('view default isometric')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('erase pid box {}'.format(erase_box))
                utils.MetaCommand('view default left')
                utils.MetaCommand('grstyle scalarfringe enable')
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('options fringebar off')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"f28_front_door".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Image 4":
                data = self.metadb_3d_input.critical_sections["f28_rear_door"]
                prop_names = data["hes"]
                erase_box = data["erase_box"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    utils.MetaCommand('add all')
                    utils.MetaCommand('view default isometric')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('erase pid box {}'.format(erase_box))
                utils.MetaCommand('view default left')
                utils.MetaCommand('grstyle scalarfringe enable')
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('options fringebar off')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"f28_front_door".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Image 2":
                data = self.metadb_3d_input.critical_sections["f21_upb_inner"]
                prop_names = data["hes"]
                erase_box = data["erase_box"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    utils.MetaCommand('add all')
                    utils.MetaCommand('view default isometric')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                self.metadb_3d_input.hide_all()
                self.metadb_3d_input.show_only_props(entities)
                utils.MetaCommand('erase pid box {}'.format(erase_box))
                utils.MetaCommand('plane new DEFAULT_PLANE_YZ xyz 1657.996826,-16.504395,576.072754 1,0,0')
                utils.MetaCommand('plane edit perpendicular 0/1/0 DEFAULT_PLANE_YZ')
                utils.MetaCommand('plane toggleopts enable sectionclip DEFAULT_PLANE_YZ')
                utils.MetaCommand('plane toggleopts disable stateauto DEFAULT_PLANE_YZ')
                utils.MetaCommand('plane options cut autovisible DEFAULT_PLANE_YZ')
                utils.MetaCommand('plane options onlysection enable DEFAULT_PLANE_YZ')
                utils.MetaCommand('plane toggleopts enable clip DEFAULT_PLANE_YZ')
                utils.MetaCommand('plane toggleopts disable sectionclip DEFAULT_PLANE_YZ')
                utils.MetaCommand('plane toggleopts enable slice DEFAULT_PLANE_YZ')
                utils.MetaCommand('plane options slicewidth 500.000000 DEFAULT_PLANE_YZ')
                utils.MetaCommand('grstyle scalarfringe enable')
                utils.MetaCommand('view default front')
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('options fringebar off')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"f21_upb_inner".lower()+".png")
                self.capture_image("MetaPost",shape.width,shape.height,image_path)
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Image 5":
                temporary_window_name = "Temporary"
                front_shoulder_intrusion_curve_name = self.general_input.front_shoulder_intrusion_curve_name
                front_door_accel_window_name = self.general_input.front_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                curve = plot2d.CurvesByName(front_door_accel_window_name, front_shoulder_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(front_door_accel_window_name,curve,temporary_window_name,"ROW 1 SHOULDER")
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(temporary_window_name, 0))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 1 SHOULDER".lower()+".png")
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 6":
                temporary_window_name = "Temporary"
                front_abdomen_intrusion_curve_name = self.general_input.front_abdomen_intrusion_curve_name
                front_door_accel_window_name = self.general_input.front_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                curve = plot2d.CurvesByName(front_door_accel_window_name, front_abdomen_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(front_door_accel_window_name,curve,temporary_window_name,"ROW 1 ABDOMEN")
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(temporary_window_name, 0))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 1 ABDOMEN".lower()+".png")
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 8":
                temporary_window_name = "Temporary"
                front_femur_intrusion_curve_name = self.general_input.front_femur_intrusion_curve_name
                front_door_accel_window_name = self.general_input.front_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                curve = plot2d.CurvesByName(front_door_accel_window_name, front_femur_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(front_door_accel_window_name,curve,temporary_window_name,"ROW 1 FEMUR")
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(temporary_window_name, 0))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 1 FEMUR".lower()+".png")
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 7":
                temporary_window_name = "Temporary"
                front_pelvis_intrusion_curve_name = self.general_input.front_pelvis_intrusion_curve_name
                front_door_accel_window_name = self.general_input.front_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                curve = plot2d.CurvesByName(front_door_accel_window_name, front_pelvis_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(front_door_accel_window_name,curve,temporary_window_name,"ROW 1 PELVIS")
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(temporary_window_name, 0))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 1 PELVIS".lower()+".png")
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 0":
                temporary_window_name = "Temporary"
                rear_shoulder_intrusion_curve_name = self.general_input.rear_shoulder_intrusion_curve_name
                rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                curve = plot2d.CurvesByName(rear_door_accel_window_name, rear_shoulder_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(rear_door_accel_window_name,curve,temporary_window_name,"ROW 2 SHOULDER")
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(temporary_window_name, 0))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 2 SHOULDER".lower()+".png")
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 10":
                temporary_window_name = "Temporary"
                rear_abdomen_intrusion_curve_name = self.general_input.rear_abdomen_intrusion_curve_name
                rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                curve = plot2d.CurvesByName(rear_door_accel_window_name, rear_abdomen_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(rear_door_accel_window_name,curve,temporary_window_name,"ROW 2 ABDOMEN")
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(temporary_window_name, 0))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 2 ABDOMEN".lower()+".png")
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 12":
                temporary_window_name = "Temporary"
                rear_femur_intrusion_curve_name = self.general_input.rear_femur_intrusion_curve_name
                rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                curve = plot2d.CurvesByName(rear_door_accel_window_name, rear_femur_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(rear_door_accel_window_name,curve,temporary_window_name,"ROW 2 FEMUR")
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(temporary_window_name, 0))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 2 FEMUR".lower()+".png")
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 11":
                temporary_window_name = "Temporary"
                rear_pelvis_intrusion_curve_name = self.general_input.rear_pelvis_intrusion_curve_name
                rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                curve = plot2d.CurvesByName(rear_door_accel_window_name, rear_pelvis_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(rear_door_accel_window_name,curve,temporary_window_name,"ROW 2 PELVIS")
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(temporary_window_name, 0))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 2 PELVIS".lower()+".png")
                img.save(image_path, 'JPEG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Table 2":
                table = shape.table
                row_index = 2
                node_index = 0
                self.report_composer.add_row(table)
                for iindex,(key,values) in enumerate(self.intrusion_areas["ROW 1"].items()):
                    if iindex == 2:
                        self.report_composer.add_row(table)
                        row_index = row_index+1
                    if iindex == 1 or iindex == 3:
                        node_index = 7
                    else:
                        node_index = 0
                    rows = table.rows
                    text_frame_1 = rows[row_index].cells[node_index].text_frame
                    font = text_frame_1.paragraphs[0].font
                    font.bold = True
                    font.size = Pt(11)
                    text_frame_1.paragraphs[0].text = key.capitalize()
                    for index,value in enumerate(values):
                        text_frame = rows[row_index].cells[node_index+index+1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(value)
            elif shape.name == "Table 3":
                table = shape.table
                row_index = 2
                node_index = 0
                self.report_composer.add_row(table)
                for iindex,(key,values) in enumerate(self.intrusion_areas["ROW 2"].items()):
                    if iindex == 2:
                        self.report_composer.add_row(table)
                        row_index = row_index+1
                    if iindex == 1 or iindex == 3:
                        node_index = 7
                    else:
                        node_index = 0
                    rows = table.rows
                    text_frame_1 = rows[row_index].cells[node_index].text_frame
                    font = text_frame_1.paragraphs[0].font
                    font.bold = True
                    font.size = Pt(11)
                    text_frame_1.paragraphs[0].text = key.capitalize()
                    for index,value in enumerate(values):
                        text_frame = rows[row_index].cells[node_index+index+1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(value)


        return 0

    def intrusion_curve_format(self,source_window,curve,temporary_window_name,curve_name):
        """
        intrusion_curve_format _summary_

        _extended_summary_

        Args:
            source_window (_type_): _description_
            curve (_type_): _description_
            temporary_window_name (_type_): _description_
            curve_name (_type_): _description_

        Returns:
            _type_: _description_
        """

        utils.MetaCommand('xyplot curve copy "{}" {}'.format(source_window,curve.id))
        utils.MetaCommand('xyplot create "{}"'.format(temporary_window_name))
        utils.MetaCommand('xyplot curve paste "{}" 0 {}'.format(temporary_window_name,curve.id))
        win = windows.Window(temporary_window_name, page_id=0)
        curve = win.get_curves('all')[0]
        y_values = []
        for x in [0.03,0.04,0.05,0.06,0.07,0.08]:
            y_values.append(round(curve.get_y_values_from_x(specifier = 'first', xvalue =x)[0]))
        self.intrusion_areas[curve_name.rsplit(" ",1)[0]][curve_name.rsplit(" ",1)[1]] = y_values
        utils.MetaCommand('xyplot gridoptions line major style "{}" 0 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions yaxis active "{}" 0 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions labels yposition "{}" 0 left'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions labels yalign "{}" 0 left'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions axyrange "{}" 0 0 0 1200'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions axyrange "{}" 0 0 0 200'.format(temporary_window_name))
        utils.MetaCommand('xyplot gridoptions yspace "{}" 0 40'.format(temporary_window_name))
        utils.MetaCommand('xyplot plotoptions title set "{}" 0 "{}"'.format(temporary_window_name,curve_name))
        utils.MetaCommand('xyplot axisoptions ylabel set "{}" 0 "Intrusion [mm]"'.format(temporary_window_name))
        utils.MetaCommand('xyplot curve select "{}" all'.format(temporary_window_name))
        utils.MetaCommand('xyplot curve set style "{}" selected 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot curve set linewidth "{}" selected 9'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions ylabel font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions labels yfont "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" 0 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions xaxis active "{}" 0 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions xlabel font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions labels xfont "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot plotoptions title font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions xaxis deactive "{}" 0 0'.format(temporary_window_name))

        return 0

    def edit_title_slide(self,slide):
        from datetime import datetime
        from pptx.util import Pt

        for shape in slide.shapes:
            if shape.name == "Table 2":
                table_obj = shape.table
                rows = table_obj.rows
                ordinal = lambda n: "%d%s" % (n,"tsnrhtdd"[(n//10%10!=1)*(n%10<4)*n%10::4])
                format_data = "%B %-d, %Y"
                month = datetime.today().strftime('%B')
                day = int(datetime.today().strftime('%d'))
                year = datetime.today().strftime('%Y')
                text_frame_1 = rows[3].cells[0].text_frame
                font = text_frame_1.paragraphs[0].font
                font.size = Pt(16)
                text_frame_2 = rows[2].cells[0].text_frame
                font = text_frame_2.paragraphs[0].font
                font.size = Pt(16)
                text_frame_3 = rows[4].cells[2].text_frame
                font = text_frame_3.paragraphs[0].font
                font.size = Pt(16)
                text_frame_1.paragraphs[0].text = " {} {}, {}".format(month,ordinal(day),year)
                text_frame_2.paragraphs[0].text = " " + self.general_input.verification_mode
                text_frame_3.paragraphs[0].text = " " + self.general_input.run_directory

        return 0

    def edit_cae_quality_slide(self,slide):

        from PIL import ImageGrab
        from pptx.util import Pt

        window_name = self.general_input.cae_quality_window_name
        utils.MetaCommand('window maximize "{}"'.format(window_name))

        for shape in slide.shapes:

            if shape.name == "Image 2":
                plot_id = 0
                page_id=0
                plot = plot2d.Plot(plot_id, window_name, page_id)
                title = plot.get_title()
                plot.activate()
                utils.MetaCommand('xyplot rlayout "{}" 1'.format(window_name))
                utils.MetaCommand('xyplot curve select "{}" all'.format(window_name))
                utils.MetaCommand('xyplot curve visible and "{}" sel'.format(window_name))
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(window_name, plot.id))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+title.get_text().lower()+".png")
                img.save(image_path, 'PNG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('xyplot plotoptions legend on "CAE quality" 0')
                utils.MetaCommand('xyplot legend hook left "CAE quality" 0')
                utils.MetaCommand('xyplot legend hook hout "CAE quality" 0')
                utils.MetaCommand('xyplot legend ymove "CAE quality" 0 1.060000')
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(window_name, plot.id))
                img_2 = ImageGrab.grabclipboard()
                legend = plot2d.Legend(plot_id, window_name, page_id)
                left,top = legend.get_position()
                width = legend.get_width()
                height = legend.get_height()
                img_2 = img_2.crop((left,top,width+8,height+8))
                image2_path = os.path.join(self.twod_images_report_folder,window_name+"_"+title.get_text().lower()+"_Legend"+".png")
                img_2.save(image2_path,"PNG")
                shape2 = [shape for shape in slide.shapes if shape.name == "Image 1"][0]
                picture = slide.shapes.add_picture(image2_path,shape2.left,shape2.top,width = shape2.width,height = shape2.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('xyplot rlayout "{}" 2'.format(window_name))
                plot.deactivate()
            elif shape.name == "Image 3":
                plot_id = 1
                page_id=0
                plot = plot2d.Plot(plot_id, window_name, page_id)
                title = plot.get_title()
                plot.activate()
                utils.MetaCommand('xyplot rlayout "{}" 1'.format(window_name))
                utils.MetaCommand('xyplot curve select "{}" all'.format(window_name))
                utils.MetaCommand('xyplot curve visible and "{}" sel'.format(window_name))
                utils.MetaCommand('clipboard copy plot image "{}" {}'.format(window_name, plot.id))
                img = ImageGrab.grabclipboard()
                img = img.resize((round(shape.width/9525),round(shape.height/9525)))
                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+title.get_text().lower()+".png")
                img.save(image_path, 'PNG')
                picture = slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                plot.deactivate()
                utils.MetaCommand('xyplot rlayout "{}" 2'.format(window_name))
            elif shape.name == "Table 1":
                plot_id = 0
                page_id=0
                plot = plot2d.Plot(plot_id, window_name, page_id)
                curvelist = plot.get_curves('all')
                index = 0
                for curve in curvelist:
                    if curve.id == 5:
                        continue
                    index = index+1
                    min_y = curve.get_limit_value_y(specifier = 'min')
                    max_y = curve.get_limit_value_y(specifier = 'max')
                    table_obj = shape.table
                    self.report_composer.add_row(table_obj)
                    row = table_obj.rows[index]
                    text_frame_1 = row.cells[0].text_frame
                    font_1 = text_frame_1.paragraphs[0].font
                    font_1.size = Pt(12)
                    text_frame_1.paragraphs[0].text = str(curve.name).replace(" energy","")

                    text_frame_2 = row.cells[1].text_frame
                    font_2 = text_frame_2.paragraphs[0].font
                    font_2.size = Pt(12)
                    text_frame_2.paragraphs[0].text = "{:.2e}".format(max_y)

                    text_frame_3 = row.cells[2].text_frame
                    font_3 = text_frame_3.paragraphs[0].font
                    font_3.size = Pt(12)
                    text_frame_3.paragraphs[0].text = "{:.2e}".format(min_y)
            elif shape.name == "Table 2":
                table_obj = shape.table
                table_value_dict ={"Termination type":self.general_input.termination_type,
                                    "Computation time":self.general_input.computation_time,
                                    "Core count":self.general_input.core_count,
                                    "Verification mode":self.general_input.verification_mode,
                                    "Compute cluster":self.general_input.compute_cluster}
                index = 0
                for item,value in table_value_dict.items():
                    index = index+1
                    self.report_composer.add_row(table_obj)
                    row = table_obj.rows[index]
                    text_frame_1 = row.cells[0].text_frame
                    font_1 = text_frame_1.paragraphs[0].font
                    font_1.size = Pt(12)
                    text_frame_1.paragraphs[0].text = item

                    if item == "Core count":
                        value = value.split("with")[1].rstrip()
                    text_frame_2 = row.cells[1].text_frame
                    font_2 = text_frame_2.paragraphs[0].font
                    font_2.size = Pt(12)
                    text_frame_2.paragraphs[0].text = value

        return 0

    def threed_data_reporting(self):
        """
        threed_data_reporting [summary]

        [extended_summary]

        Returns:
            [type]: [description]
        """
        critical_sections_data = self.metadb_3d_input.critical_sections
        # for section,value in critical_sections_data.items():
        #     for key,vvalue in value.items():
        #         if key == "hes":

        return 0

    def twod_data_reporting(self):
        """
        twod_data_reporting [summary]

        [extended_summary]

        Returns:
            [type]: [description]
        """
        from PIL import ImageGrab

        window_2d_objects = self.metadb_2d_input.window_objects
        for window in window_2d_objects:
            window_name = window.name
            window_layout = window.meta_obj.get_plot_layout()
            plot = window.plot
            curve = plot.curve
            utils.MetaCommand('window active "{}"'.format(window_name))
            utils.MetaCommand('window maximize "{}"'.format(window_name))
            utils.MetaCommand('xyplot plotdeactive "{}" all'.format(window_name))
            curve.meta_obj.show()
            utils.MetaCommand('xyplot plotactive "{}" {}'.format(window_name, plot.id))
            utils.MetaCommand('xyplot curve visible and "{}" selected'.format(window_name))
            #utils.MetaCommand('xyplot rlayout "{}" 1'.format(window_name))
            image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+curve.name.lower()+".png")
            if not os.path.exists(os.path.dirname(image_path)):
                print(os.path.dirname(image_path))
                os.makedirs(os.path.dirname(image_path))
            utils.MetaCommand('clipboard copy plot image "{}" {}'.format(window_name, plot.id))
            img = ImageGrab.grabclipboard()
            img.save(image_path, 'PNG')
            #utils.MetaCommand('write jpeg "{}" 100'.format(image_path))
            #utils.MetaCommand('xyplot rlayout "{}" {}'.format(window_name,window_layout))

        return 0


class PPTXReportComposer():
    def __init__(self, report_name, template_pptx):

        self.report_name = report_name
        self.template_pptx = template_pptx
        self.prs_obj = None

    def create_prs_obj(self):
        """ Creates the PPTx report using
        the python-pptx module
        """
        from pptx import Presentation

        # Instantiate
        if not self.prs_obj:
            self.prs_obj = Presentation(self.template_pptx)

        return 0

    # def add_slide(self,layout_name):

    #     slide = self.prs_obj.slides.add_slide(layout_name)

    #     return slide

    @staticmethod
    def add_row(table):
        """
        add_row [summary]

        [extended_summary]

        Args:
            table ([type]): [description]

        Returns:
            [type]: [description]
        """
        from pptx.table import _Cell
        from copy import deepcopy

        new_row = deepcopy(table._tbl.tr_lst[-1])
        # duplicating last row of the table as a new row to be added

        table._tbl.append(new_row)

        return 0

    @staticmethod
    def remove_row(table,row_to_delete):
        """
        remove_row [summary]

        [extended_summary]

        Args:
            table ([type]): [description]
            row_to_delete ([type]): [description]
        """
        table._tbl.remove(row_to_delete._tr)

    def save_pptx(self, pptx_filepath, datestamp=""):
        """ Saves the PPTx at the given filepath
        with the given datestamp

        Args:
            pptx_filepath (str): Absolute path to the pptx file for saving.
            datestamp (str, optional): Date stamp at the bottom right of slide. Defaults to "".

        Returns:
            int: 0 always.
        """
        from pptx.enum.text import PP_ALIGN

        # Get current date if not provided
        if not datestamp:
            datestamp = time.strftime("%B %d, %Y")

        # Set the date in the ppt master slide
        master_slide = self.prs_obj.slide_master
        for shape in master_slide.shapes:
            try:
                if shape.text == "Date_Stamp":
                    shape.text = datestamp
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            except AttributeError:
                continue

        # Save the pptx
        self.prs_obj.save(pptx_filepath)

        return 0
