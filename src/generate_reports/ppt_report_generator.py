# PYTHON script
"""
This Script is used to generate all the slides of the Thesis Report Generation.
"""

import os
import time
import logging

from src.generate_reports.report_slides.title_slide import TitleSlide
from src.generate_reports.report_slides.cae_quality_slide import CAEQualitySlide
from src.generate_reports.report_slides.executive_slide import ExecutiveSlide
from src.generate_reports.report_slides.cbu_and_barrier_position_slide import CBUAndBarrierPositionSlide
from src.generate_reports.report_slides.biw_kinematics_slide import BIWKinematicsSlide
from src.generate_reports.report_slides.biw_cbu_deformation_slide import BIWCBUDeformationSlide
from src.generate_reports.report_slides.bom_f21_upb_slide import BOMF21UPBSlide
from src.generate_reports.report_slides.biw_stiff_ring_deformation_slide import BIWStiffRingDeformationSlide
from src.generate_reports.report_slides.biw_bplr_def_and_intr_slide import BIWBplrDeformationAndIntrusion
from src.generate_reports.report_slides.bom_f21_roof_slide import BOMF21ROOFSlide
from src.generate_reports.report_slides.biw_roof_def_and_spot_failure import BIWROOFDeformationAndSpotWeldFailure
from src.generate_reports.report_slides.bom_f28_doors_slide import BOMF28DoorsSlide
from src.generate_reports.report_slides.bom_f21_front_floor import BOMF21FrontFloorSlide
from src.generate_reports.report_slides.biw_floor_deformation_spotweld_failure_slide import BIWFloorDeformationAndSpotWeldFailureSlide
from src.generate_reports.report_slides.enclosure_performance_skin_deformation_slide import EnclosurePerformanceSkinDeformationSlide
from src.generate_reports.report_slides.bom_row2_f28_doors_slide import BOMRow2F28DoorsSlide
from src.generate_reports.report_slides.enclosure_performance_front_door_panel_intrusion_slide import EnclosurePerformanceFrontDoorPanelIntrusionSlide
from src.generate_reports.report_slides.enclosure_performance_front_door_panel_deformation_slide import EnclosurePerformanceFrontDoorPanelDeformationSlide
from src.generate_reports.report_slides.enclosure_performance_rear_door_panel_deformation_slide import EnclosurePerformanceRearDoorPanelDeformationSlide
from src.generate_reports.report_slides.enclosures_performance_rear_door_panel_intrusion_slide import EnclosurePerformanceRearDoorPanelIntrusionSlide
from src.generate_reports.report_slides.biw_stiff_ring_spotweld_failure_slide import BIWStiffRingSpotWeldFailureSlide

class SideCrashPPTReportGenerator():
    """
        This Class is used to generate the slides of thesis Report.

        Args:
            windows (_type_): _description_
            general_input (GeneralInfo): GeneralInfo class object.
            metadb_2d_input (Meta2DInfo): Meta2DInfo class object.
            metadb_3d_input (Meta3DInfo): Meta3DInfo class object.
            template_file (_type_): _description_
            twod_images_report_folder (str): folder path to save twod data images.
            threed_images_report_folder (str): folder path to save threed data images.
            ppt_report_folder (str): folder path to save the ppt
    """

    def __init__(self,
                windows,
                general_input,
                metadb_2d_input,
                metadb_3d_input,
                template_file,
                twod_images_report_folder,
                threed_images_report_folder,
                ppt_report_folder) -> None:

        self.windows = windows
        self.general_input = general_input
        self.metadb_2d_input = metadb_2d_input
        self.metadb_3d_input = metadb_3d_input
        self.template_file = template_file
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.ppt_report_folder = ppt_report_folder
        self.report_composer = None
        self.logger = logging.getLogger("side_crash_logger")

    def generate_ppt(self):
        """
        This method is used to generate the Slides of thesis Report

        Returns:
            [0]: 0 for Success 1 for failure
        """
        #pptx object creation with template
        self.report_composer = PPTXReportComposer(report_name="Run1",template_pptx=self.template_file)
        self.report_composer.create_prs_obj()
        #editing slides of the created pptx object
        title_slide = TitleSlide(self.report_composer.prs_obj.slides[0],
                                self.general_input)
        title_slide.edit()
        cae_quality_slide = CAEQualitySlide(self.report_composer.prs_obj.slides[1],
                                self.general_input,
                                self.twod_images_report_folder)
        cae_quality_slide.edit()
        executive_slide = ExecutiveSlide(self.report_composer.prs_obj.slides[2],
                                self.general_input,
                                self.metadb_3d_input,
                                self.twod_images_report_folder,
                                self.threed_images_report_folder)
        executive_slide.edit()
        # cbu_and_barrier_position_slide = CBUAndBarrierPositionSlide(self.report_composer.prs_obj.slides[3],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.threed_images_report_folder)
        # cbu_and_barrier_position_slide.edit()
        # biw_kinematics_slide = BIWKinematicsSlide(self.report_composer.prs_obj.slides[6],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.twod_images_report_folder,
        #                         self.threed_images_report_folder)
        # biw_kinematics_slide.edit()
        # biw_cbu_deformation_slide = BIWCBUDeformationSlide(self.report_composer.prs_obj.slides[7],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.threed_images_report_folder)
        # biw_cbu_deformation_slide.edit()
        # bom_f21_upb_slide = BOMF21UPBSlide(self.report_composer.prs_obj.slides[8],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.threed_images_report_folder)
        # bom_f21_upb_slide.edit()
        # biw_stiff_ring_deformation_slide = BIWStiffRingDeformationSlide(self.report_composer.prs_obj.slides[9],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.twod_images_report_folder,
        #                         self.threed_images_report_folder)
        # biw_stiff_ring_deformation_slide.edit()
        # biw_bplr_def_and_intr_slide = BIWBplrDeformationAndIntrusion(self.report_composer.prs_obj.slides[10],
        #                          self.general_input,
        #                         self.metadb_3d_input,
        #                         self.twod_images_report_folder,
        #                         self.threed_images_report_folder)
        # biw_bplr_def_and_intr_slide.edit()
        # bom_f21_roof_slide = BOMF21ROOFSlide(self.report_composer.prs_obj.slides[12],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.threed_images_report_folder)
        # bom_f21_roof_slide.edit()
        # biw_roof_def_and_spot_failure_slide = BIWROOFDeformationAndSpotWeldFailure(self.report_composer.prs_obj.slides[13],
        #                             self.general_input,
        #                             self.metadb_3d_input,
        #                             self.twod_images_report_folder,
        #                             self.threed_images_report_folder)
        # biw_roof_def_and_spot_failure_slide.edit()
        # bom_row2_f28_doors_slide = BOMRow2F28DoorsSlide(self.report_composer.prs_obj.slides[25],
        #                             self.general_input,
        #                             self.metadb_3d_input,
        #                             self.threed_images_report_folder)
        # bom_row2_f28_doors_slide.edit()
        # bom_f21_front_floor_slide = BOMF21FrontFloorSlide(self.report_composer.prs_obj.slides[15],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.threed_images_report_folder)
        # bom_f21_front_floor_slide.edit()
        # biw_floor_deformation_and_spotweld_failure_slide = BIWFloorDeformationAndSpotWeldFailureSlide(self.report_composer.prs_obj.slides[16],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.twod_images_report_folder,
        #                         self.threed_images_report_folder)
        # biw_floor_deformation_and_spotweld_failure_slide.edit()
        # enclosure_performance_skin_deformation_slide = EnclosurePerformanceSkinDeformationSlide(self.report_composer.prs_obj.slides[17],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.twod_images_report_folder,
        #                         self.threed_images_report_folder)
        # enclosure_performance_skin_deformation_slide.edit()
        # bom_f28_doors_slide = BOMF28DoorsSlide(self.report_composer.prs_obj.slides[18],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.threed_images_report_folder)
        # bom_f28_doors_slide.edit()
        # enclosure_performance_front_door_panel_intrusion_slide = EnclosurePerformanceFrontDoorPanelIntrusionSlide(self.report_composer.prs_obj.slides[19],
        #                         self.general_input,
        #                         self.twod_images_report_folder)
        # enclosure_performance_front_door_panel_intrusion_slide.edit()
        # enclosure_performance_front_door_panel_deformation_slide = EnclosurePerformanceFrontDoorPanelDeformationSlide(self.report_composer.prs_obj.slides[20],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.twod_images_report_folder,
        #                         self.threed_images_report_folder)
        # enclosure_performance_front_door_panel_deformation_slide.edit()
        # enclosure_performance_rear_door_panel_deformation_slide = EnclosurePerformanceRearDoorPanelDeformationSlide(self.report_composer.prs_obj.slides[27],
        #                         self.general_input,
        #                         self.metadb_3d_input,
        #                         self.twod_images_report_folder,
        #                         self.threed_images_report_folder)
        # enclosure_performance_rear_door_panel_deformation_slide.edit()
        # enclosures_performance_rear_door_panel_intrusion_slide = EnclosurePerformanceRearDoorPanelIntrusionSlide(self.report_composer.prs_obj.slides[26],
        #                         self.general_input,
        #                         self.twod_images_report_folder)
        # enclosures_performance_rear_door_panel_intrusion_slide.edit()
        biw_stiff_ring_spotweld_failure_slide = BIWStiffRingSpotWeldFailureSlide(self.report_composer.prs_obj.slides[11],
                                self.general_input,
                                self.metadb_3d_input,
                                self.threed_images_report_folder)
        biw_stiff_ring_spotweld_failure_slide.edit()
        sa
        #saving the pptx object
        file_name = os.path.join(self.ppt_report_folder,"Side_MDB_Thesis_Report_{}.pptx".format(time.strftime("%Y-%d-%m")))
        self.report_composer.save_pptx(file_name)

        return 0

class PPTXReportComposer():
    """
    This Class is used to Create the Presentation of the template and saving the thesis report and Executive Report.

    Args:
        report_name (str): name of the Report
        template_pptx (obj): template of the pptx
    """
    def __init__(self, report_name, template_pptx):
        self.report_name = report_name
        self.template_pptx = template_pptx
        self.prs_obj = None

    def create_prs_obj(self):
        """ Creates the PPTx report using
        the python-pptx module
        """
        from pptx import Presentation

        # Instantiate
        if not self.prs_obj:
            self.prs_obj = Presentation(self.template_pptx)

        return 0

    def save_pptx(self, pptx_filepath, datestamp=""):
        """ Saves the PPTx at the given filepath
        with the given datestamp

        Args:
            pptx_filepath (str): Absolute path to the pptx file for saving.
            datestamp (str, optional): Date stamp at the bottom right of slide. Defaults to "".

        Returns:
            int: 0 always.
        """
        from pptx.enum.text import PP_ALIGN
        from pptx import Presentation

        # Get current date if not provided
        if not datestamp:
            datestamp = time.strftime("%B %d, %Y")

        # Set the date in the ppt master slide
        master_slide = self.prs_obj.slide_master
        for shape in master_slide.shapes:
            try:
                if shape.text == "Date_Stamp":
                    shape.text = datestamp
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            except AttributeError:
                continue

        # Save the pptx
        self.prs_obj.save(pptx_filepath)

        # Thesis and Executive Report Presentation
        thesis_report_presentation = Presentation(pptx_filepath)
        executive_report_presentation = Presentation(pptx_filepath)
        # Thesis and Executive Report Slides
        thesis_report_xml_slides = thesis_report_presentation.slides._sldIdLst
        executive_reportxml2_slides = executive_report_presentation.slides._sldIdLst
        thesis_report_slides = list(thesis_report_xml_slides)
        executive_report_slides = list(executive_reportxml2_slides)
        # Removing the 2nd slide from Thesis Report and saving
        for index,_slide in enumerate(thesis_report_slides):
            if index != 2:
                thesis_report_xml_slides.remove(thesis_report_slides[index])
        thesis_report_presentation.save(pptx_filepath.replace("Thesis","Executive"))
        # Saving only the 2nd slide and saving
        for index,_slide2 in enumerate(executive_report_slides):
            if index == 2:
                executive_reportxml2_slides.remove(executive_report_slides[index])

        executive_report_presentation.save(pptx_filepath)

        return 0
