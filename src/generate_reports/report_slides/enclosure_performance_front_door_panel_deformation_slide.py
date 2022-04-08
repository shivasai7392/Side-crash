# PYTHON script
"""
This script is used for all the automation process of Enclosure performance front door panel deformation slide of thesis report.
"""


import os
import logging
from datetime import datetime

from meta import utils
from meta import plot2d
from meta import windows

from src.meta_utilities import capture_image
from src.meta_utilities import  visualize_3d_critical_section
from src.meta_utilities import capture_image_and_resize
from src.general_utilities import closest
from src.metadb_info import GeneralVarInfo

class EnclosurePerformanceFrontDoorPanelDeformationSlide():
    """
       This class is used to automate the enclosure performance front door panel deformation slide of thesis report.

        Args:
            slide (object): enclosure performance front door panel intrsusion pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
            metadb_3d_input (Meta3DInfo): Meta3DInfo class object.
            twod_images_report_folder (str): folder path to save twod data images.
            threed_images_report_folder (str): folder path to save threed data images.
        """
    def __init__(self,
                slide,
                general_input,
                metadb_3d_input,
                twod_images_report_folder,
                threed_images_report_folder) -> None:
        self.shapes = slide.shapes
        self.general_input = general_input
        self.metadb_3d_input = metadb_3d_input
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.logger = logging.getLogger("side_crash_logger")

    def edit(self):
        """
        This method is used to iterate all the shapes of the enclosure performance front door panel deformation slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """
        from pptx.util import Pt
        try:
            self.logger.info("Started seeding data into enclosure performance front door panel deformation slide")
            self.logger.info("")
            starttime = datetime.now()
            #iterating through the shapes of the enclosure performance front door panel deformation slide
            for shape in self.shapes:
                #image insertion for the shape named "Image 1"
                if shape.name == "Image 1":
                    #capturing fringe bar of metapost window
                    utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                    utils.MetaCommand('0:options state variable "serial=1"')
                    utils.MetaCommand('options fringebar on')
                    if self.threed_images_report_folder is not None:
                        image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"FRINGE_BAR"+".png").replace(" ","_")
                        utils.MetaCommand('write scalarfringebar png {} '.format(image_path))
                        self.logger.info("--- 3D FRINGE BAR IMAGE GENERATOR")
                        self.logger.info("")
                        self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                        self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                        self.logger.info("OUTPUT MODEL IMAGES :")
                        self.logger.info(image_path)
                        self.logger.info("")
                        #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                        picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                        picture.crop_left = 0
                        picture.crop_right = 0
                        utils.MetaCommand('options fringebar off')
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                        self.logger.info("")
                #image insertion for the shape named "Image 2"
                elif shape.name == "Image 2":
                    #visualizing and capturing "f28_front_door" critical part set at peak state with deformation
                    utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                    data = self.metadb_3d_input.critical_sections["f28_front_door"]
                    visualize_3d_critical_section(data)
                    utils.MetaCommand('0:options state variable "serial=1"')
                    utils.MetaCommand('grstyle scalarfringe enable')
                    utils.MetaCommand('options fringebar off')
                    utils.MetaCommand('grstyle deform on')
                    if self.threed_images_report_folder is not None:
                        image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"F28_FRONT_DOOR_AT_PEAK_STATE_WITH_DEFORMATION"+".png").replace(" ","_")
                        capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height,view = "left")
                        self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                        self.logger.info("")
                        self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                        self.logger.info("SOURCE MODEL : 0")
                        self.logger.info("STATE : PEAK STATE")
                        self.logger.info("PID NAME SHOW FILTER : {} ".format(data["hes"] if "hes" in data.keys() else "null"))
                        self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                        self.logger.info("PID NAME ERASE FILTER : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                        self.logger.info("PID'S TO ERASE : {} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null"))
                        self.logger.info("ERASE BOX : {} ".format(data["erase_box"] if "erase_box" in data.keys() else "null"))
                        self.logger.info("IMAGE VIEW : {} ".format(data["view"] if "view" in data.keys() else "null"))
                        self.logger.info("TRANSPARENCY LEVEL : 50" )
                        self.logger.info("TRANSPARENT PID'S : {} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null"))
                        self.logger.info("COMP NAME : {} ".format(data["name"] if "name" in data.keys() else "null"))
                        self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                        self.logger.info("OUTPUT MODEL IMAGES :")
                        self.logger.info(image_path)
                        self.logger.info("")
                        #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                        picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                        picture.crop_left = 0
                        picture.crop_right = 0
                        #reverting visual settings
                        utils.MetaCommand('color pid transparency reset act')
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                        self.logger.info("")
                #image insertion for the shape named "Image 3"
                elif shape.name == "Image 3":
                    #visualizing and capturing "f28_front_door" critical part set at peak state without deformation
                    utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                    data = self.metadb_3d_input.critical_sections["f28_front_door"]
                    visualize_3d_critical_section(data)
                    utils.MetaCommand('0:options state variable "serial=1"')
                    utils.MetaCommand('grstyle scalarfringe enable')
                    utils.MetaCommand('options fringebar off')
                    utils.MetaCommand('grstyle deform off')
                    if self.threed_images_report_folder is not None:
                        image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"F28_FRONT_DOOR_AT_PEAK_STATE_WITH_DEFORMATION"+".png").replace(" ","_")
                        capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height,view = "left")
                        self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                        self.logger.info("")
                        self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                        self.logger.info("SOURCE MODEL : 0")
                        self.logger.info("STATE : PEAK STATE")
                        self.logger.info("PID NAME SHOW FILTER : {} ".format(data["hes"] if "hes" in data.keys() else "null"))
                        self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                        self.logger.info("PID NAME ERASE FILTER : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                        self.logger.info("PID'S TO ERASE : {} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null"))
                        self.logger.info("ERASE BOX : {} ".format(data["erase_box"] if "erase_box" in data.keys() else "null"))
                        self.logger.info("IMAGE VIEW : {} ".format(data["view"] if "view" in data.keys() else "null"))
                        self.logger.info("TRANSPARENCY LEVEL : 50" )
                        self.logger.info("TRANSPARENT PID'S : {} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null"))
                        self.logger.info("COMP NAME : {} ".format(data["name"] if "name" in data.keys() else "null"))
                        self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                        self.logger.info("OUTPUT MODEL IMAGES :")
                        self.logger.info(image_path)
                        self.logger.info("")
                        #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                        picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                        picture.crop_left = 0
                        picture.crop_right = 0
                        #reverting visual settings
                        utils.MetaCommand('grstyle deform on')
                        utils.MetaCommand('color pid transparency reset act')
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                        self.logger.info("")
                #image insertion for the shape named "Image 4"
                elif shape.name == "Image 4":
                    if not self.general_input.door_panel_intrusion_window_name in ["null","none",""]:
                        door_panel_intrusion_window_name = self.general_input.door_panel_intrusion_window_name
                        door_panel_intrusion_window_obj = windows.WindowByName(door_panel_intrusion_window_name)
                        if door_panel_intrusion_window_obj:
                            #getting "Door panel intrusion" window initial,final and peak time curve object of plots with ids 0,2,4
                            utils.MetaCommand('window maximize {}'.format(door_panel_intrusion_window_name))
                            door_panel_intrusion_window = windows.Window(door_panel_intrusion_window_name, page_id=0)
                            door_panel_intrusion_window_layout = door_panel_intrusion_window.get_plot_layout()
                            door_skin_intrusion_window_plots = [plot2d.Plot(plot_id, door_panel_intrusion_window_name, 0) for plot_id in [0,2,4]]
                            for plot in door_skin_intrusion_window_plots:
                                plot.activate()
                                initial_curve_re  ="*0MS"
                                curves = plot.get_curves('all')
                                deformation_line_list = []
                                if self.general_input.survival_space_final_time not in ["null","none",""]:
                                    final_time = str(round(float(self.general_input.survival_space_final_time)))
                                    final_time_curve_name = "*{}MS".format(final_time)
                                    final_curves = plot.get_curves('byname', name = final_time_curve_name)
                                    if final_curves:
                                        final_curves[0].show()
                                    else:
                                        self.logger.info("ERROR : Door panel intrusion window does not contain Front '{}' curve from META 2D variable {}. Please update.".format(final_time_curve_name,GeneralVarInfo.survival_space_final_time_key))
                                        self.logger.info("")
                                else:
                                    self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.survival_space_final_time_key))
                                    self.logger.info("")
                                for each_curve in curves:
                                    ms = each_curve.name.rsplit("_",1)[1]
                                    if 'MS' in ms:
                                        ms_replacing = ms.replace('MS',"")
                                        deformation_line_list.append(int(ms_replacing))
                                if self.general_input.peak_time_display_value not in ["null","none",""]:
                                    peak_time = str(closest(deformation_line_list, round(float(self.general_input.peak_time_display_value))))
                                    peak_curve_re = "*{}MS".format(peak_time)
                                    peak_curves = plot.get_curves('byname', name = peak_curve_re)
                                    if peak_curves:
                                        peak_curves[0].show()
                                        peak_curves[0].set_line_style(line_style = 5)
                                    else:
                                        self.logger.info("ERROR : Door panel intrusion window does not contain Front '{}' curve from META 2D variable {}. Please update.".format(peak_curve_re,GeneralVarInfo.peak_time_display_key))
                                        self.logger.info("")
                                else:
                                    self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.peak_time_display_key))
                                    self.logger.info("")
                                initial_curve = plot.get_curves('byname', name = initial_curve_re)[0]
                                initial_curve.show()
                                utils.MetaCommand('xyplot axisoptions yaxis active "{}" {} 0'.format(door_panel_intrusion_window_name,plot.id))
                                utils.MetaCommand('xyplot axisoptions ylabel font "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(door_panel_intrusion_window_name,plot.id))
                                utils.MetaCommand('xyplot axisoptions labels yfont "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(door_panel_intrusion_window_name,plot.id))
                                utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" {} 0'.format(door_panel_intrusion_window_name,plot.id))
                                utils.MetaCommand('xyplot axisoptions xaxis active "{}" {} 0'.format(door_panel_intrusion_window_name,plot.id))
                                utils.MetaCommand('xyplot axisoptions xlabel font "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(door_panel_intrusion_window_name,plot.id))
                                utils.MetaCommand('xyplot axisoptions labels xfont "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(door_panel_intrusion_window_name,plot.id))
                                utils.MetaCommand('xyplot axisoptions xaxis deactive "{}" {} 0'.format(door_panel_intrusion_window_name,plot.id))
                                utils.MetaCommand('xyplot plotoptions title font "{}" {} "Arial,12,-1,5,75,0,0,0,0,0"'.format(door_panel_intrusion_window_name,plot.id))
                            utils.MetaCommand('xyplot rlayout "Door panel intrusion" 6')
                            if self.twod_images_report_folder is not None:
                                #capturing image of "Door panel intrusion window"
                                image_path = os.path.join(self.twod_images_report_folder,door_panel_intrusion_window_name+"FRONT_DOOR_CURVES"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | SOURCE PLOT : VISIBLE PLOTS WITH IDS 0,2,4 | SOURCE WINDOW : {}".format("*0MS,*{}MS,*{}MS".format(final_time,peak_time),door_panel_intrusion_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #reverting window plot settings and attribues
                                utils.MetaCommand('xyplot rlayout "Door panel intrusion" {}'.format(door_panel_intrusion_window_layout))
                                utils.MetaCommand('xyplot plotdeactive "{}" all'.format(door_panel_intrusion_window_name))
                            else:
                                self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        else:
                            self.logger.info("ERROR : 2D METADB does not contain 'Door panel intrusion'. Please update.")
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.front_door_accel_window_key))
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into enclosure performance front door panel deformation slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into enclosure performance front door panel deformation slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")

        return 0
