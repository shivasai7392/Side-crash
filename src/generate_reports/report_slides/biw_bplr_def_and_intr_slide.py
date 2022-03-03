# PYTHON script
"""
    _summary_

_extended_summary_

Returns:
    _type_: _description_
"""
import os

from meta import utils
from meta import plot2d

from src.meta_utilities import capture_image,visualize_3d_critical_section
from src.meta_utilities import capture_resized_image
from src.general_utilities import add_row
from src.general_utilities import closest

class BIWBplrDeformationAndIntrusion():

    def __init__(self,
                slide,
                meta_windows,
                general_input,
                metadb_2d_input,
                metadb_3d_input,
                template_file,
                twod_images_report_folder,
                threed_images_report_folder,
                ppt_report_folder) -> None:
        self.shapes = slide.shapes
        self.meta_windows = meta_windows
        self.general_input = general_input
        self.meta2d_input = metadb_2d_input
        self.meta3d_input = metadb_3d_input
        self.template_file = template_file
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.ppt_report_folder = ppt_report_folder


    def edit(self):
        """
        edit _summary_

        _extended_summary_
        """
        from pptx.util import Pt
        window_name = self.general_input.survival_space_window_name
        cds
        utils.MetaCommand('window maximize "{}"'.format(window_name))
        for shape in self.shapes:
            if shape.name == "Image 1":
                final_time_curve_name = "SS_{}MS".format(round(float(self.general_input.survival_space_final_time)))
                initial_time_curve_name = "SS_0MS"
                plot_id = 0
                page_id=0
                plot = plot2d.Plot(plot_id, window_name, page_id)
                title = plot.get_title()
                org_name = title.get_text()
                org_name = title.get_text()
                title.set_text("{} 0MS AND 150MS".format(org_name))
                plot.activate()
                final_curve = plot2d.CurvesByName(window_name, final_time_curve_name, 0)[0]
                final_curve.show()
                initial_curve = plot2d.CurvesByName(window_name, initial_time_curve_name, 0)[0]
                initial_curve.show()
                peak_time_value = self.general_input.peak_time_display_value
                peak_time_value = peak_time_value.split(".")[0]

                survival_space_curves = plot.get_curves('all')
                survival_space_cuves_list = list()
                for each_survival_space_curves in survival_space_curves:
                    ms = each_survival_space_curves.name.split("_")[1]
                    if 'MS' in ms:
                        ms_replacing = ms.replace('MS',"")
                        survival_space_cuves_list.append(int(ms_replacing))
                peak_curve_value = closest(survival_space_cuves_list, int(peak_time_value))
                peak_curve_value = plot.get_curves('byname', name ="SS_"+str(peak_curve_value)+"MS")
                peak_curve = plot2d.CurvesByName(window_name, peak_curve_value[0].name, 0)[0]
                peak_curve.show()
                utils.MetaCommand('xyplot curve set color "{}" {} "Cyan"'.format(window_name, initial_curve.id))
                utils.MetaCommand('xyplot curve set style "{}" {} 9'.format(window_name, initial_curve.id))
                utils.MetaCommand('xyplot curve set style "{}" {} 5'.format(window_name, peak_curve.id))

                image_path = os.path.join(self.twod_images_report_folder,window_name+"_with_peak_time"+title.get_text().lower()+".png")
                capture_image(window_name,shape.width,shape.height,image_path,plot_id=plot.id)
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                title.set_text(org_name)
                picture.crop_left = 0
                picture.crop_right = 0
