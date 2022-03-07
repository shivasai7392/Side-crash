# PYTHON script
"""
    _summary_

_extended_summary_

Returns:
    _type_: _description_
"""

import os

from meta import utils
from meta import plot2d
from meta import windows

from src.meta_utilities import capture_resized_image
from src.general_utilities import add_row

class EnclosurePerformanceFrontDoorPanelIntrusionSlide():

    def __init__(self,
                slide,
                meta_windows,
                general_input,
                metadb_2d_input,
                metadb_3d_input,
                template_file,
                twod_images_report_folder,
                threed_images_report_folder,
                ppt_report_folder) -> None:
        self.shapes = slide.shapes
        self.windows = meta_windows
        self.general_input = general_input
        self.metadb_2d_input = metadb_2d_input
        self.metadb_3d_input = metadb_3d_input
        self.template_file = template_file
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.ppt_report_folder = ppt_report_folder
        self.visible_parts = None
        self.intrusion_areas = {"ROW 1":{},"ROW 2":{}}

    def intrusion_curve_format(self,source_window,curve,temporary_window_name,curve_name):
        """
        intrusion_curve_format _summary_

        _extended_summary_

        Args:
            source_window (_type_): _description_
            curve (_type_): _description_
            temporary_window_name (_type_): _description_
            curve_name (_type_): _description_

        Returns:
            _type_: _description_
        """

        utils.MetaCommand('xyplot curve copy "{}" {}'.format(source_window,curve.id))
        utils.MetaCommand('xyplot create "{}"'.format(temporary_window_name))
        utils.MetaCommand('xyplot curve paste "{}" 0 {}'.format(temporary_window_name,curve.id))
        win = windows.Window(temporary_window_name, page_id=0)
        curve = win.get_curves('all')[0]
        y_values = []
        for x in [0.03,0.04,0.05,0.06,0.07,0.08]:
            y_values.append(round(curve.get_y_values_from_x(specifier = 'first', xvalue =x)[0]))
        self.intrusion_areas[curve_name.rsplit(" ",1)[0]][curve_name.rsplit(" ",1)[1]] = y_values
        utils.MetaCommand('xyplot gridoptions line major style "{}" 0 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions yaxis active "{}" 0 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions labels yposition "{}" 0 left'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions labels yalign "{}" 0 left'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions axyrange "{}" 0 0 0 1200'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions axyrange "{}" 0 0 0 200'.format(temporary_window_name))
        utils.MetaCommand('xyplot gridoptions yspace "{}" 0 40'.format(temporary_window_name))
        utils.MetaCommand('xyplot plotoptions title set "{}" 0 "{}"'.format(temporary_window_name,curve_name))
        utils.MetaCommand('xyplot axisoptions ylabel set "{}" 0 "Intrusion [mm]"'.format(temporary_window_name))
        utils.MetaCommand('xyplot curve select "{}" all'.format(temporary_window_name))
        utils.MetaCommand('xyplot curve set style "{}" selected 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot curve set linewidth "{}" selected 9'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions ylabel font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions labels yfont "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" 0 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions xaxis active "{}" 0 0'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions xlabel font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions labels xfont "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot plotoptions title font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
        utils.MetaCommand('xyplot axisoptions xaxis deactive "{}" 0 0'.format(temporary_window_name))

        return 0

    def setup(self):
        """
        setup _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """

        return 0

    def edit(self):
        """
        edit _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """
        from pptx.util import Pt

        self.setup()
        window_name = self.general_input.survival_space_window_name
        for shape in self.shapes:
            if shape.name == "Image 1":
                temporary_window_name = "Temporary"
                front_shoulder_intrusion_curve_name = self.general_input.front_shoulder_intrusion_curve_name
                front_door_accel_window_name = self.general_input.front_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                curve = plot2d.CurvesByName(front_door_accel_window_name, front_shoulder_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(front_door_accel_window_name,curve,temporary_window_name,"ROW 1 SHOULDER")

                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 1 SHOULDER".lower()+".png")
                capture_resized_image(temporary_window_name,shape.width,shape.height,image_path,plot_id=0)
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 2":
                temporary_window_name = "Temporary"
                front_abdomen_intrusion_curve_name = self.general_input.front_abdomen_intrusion_curve_name
                front_door_accel_window_name = self.general_input.front_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                curve = plot2d.CurvesByName(front_door_accel_window_name, front_abdomen_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(front_door_accel_window_name,curve,temporary_window_name,"ROW 1 ABDOMEN")

                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 1 ABDOMEN".lower()+".png")
                capture_resized_image(temporary_window_name,shape.width,shape.height,image_path,plot_id=0)
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 4":
                temporary_window_name = "Temporary"
                front_femur_intrusion_curve_name = self.general_input.front_femur_intrusion_curve_name
                front_door_accel_window_name = self.general_input.front_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                curve = plot2d.CurvesByName(front_door_accel_window_name, front_femur_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(front_door_accel_window_name,curve,temporary_window_name,"ROW 1 FEMUR")

                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 1 FEMUR".lower()+".png")
                capture_resized_image(temporary_window_name,shape.width,shape.height,image_path,plot_id=0)
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            elif shape.name == "Image 3":
                temporary_window_name = "Temporary"
                front_pelvis_intrusion_curve_name = self.general_input.front_pelvis_intrusion_curve_name
                front_door_accel_window_name = self.general_input.front_door_accel_window_name
                utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                curve = plot2d.CurvesByName(front_door_accel_window_name, front_pelvis_intrusion_curve_name, 1)[0]
                curve.show()
                self.intrusion_curve_format(front_door_accel_window_name,curve,temporary_window_name,"ROW 1 PELVIS")

                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+"ROW 1 PELVIS".lower()+".png")
                capture_resized_image(temporary_window_name,shape.width,shape.height,image_path,plot_id=0)
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
        for shape in self.shapes:
            if shape.name == "Table 1":
                table = shape.table
                row_index = 2
                node_index = 0
                add_row(table)
                for iindex,(key,values) in enumerate(self.intrusion_areas["ROW 1"].items()):
                    if iindex == 2:
                        add_row(table)
                        row_index = row_index+1
                    if iindex == 1 or iindex == 3:
                        node_index = 7
                    else:
                        node_index = 0
                    rows = table.rows
                    text_frame_1 = rows[row_index].cells[node_index].text_frame
                    font = text_frame_1.paragraphs[0].font
                    font.bold = True
                    font.size = Pt(11)
                    text_frame_1.paragraphs[0].text = key.capitalize()
                    for index,value in enumerate(values):
                        text_frame = rows[row_index].cells[node_index+index+1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(value)
        self.revert()

        return 0

    def revert(self):
        """
        revert _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """

        return 0
