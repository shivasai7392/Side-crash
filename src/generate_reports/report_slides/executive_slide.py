# PYTHON script
"""
This script is used for all the automation process of Executive report.
"""
import logging
import os
from datetime import datetime

from meta import utils
from meta import plot2d
from meta import windows

from src.meta_utilities import capture_image
from src.meta_utilities import visualize_3d_critical_section
from src.meta_utilities import capture_image_and_resize
from src.general_utilities import add_row
from src.metadb_info import GeneralVarInfo

class ExecutiveSlide():
    """
        This class is used to automate the Executive report.

        Args:
            slide (object): executive report pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
            metadb_3d_input (Meta3DInfo): Meta3DInfo class object.
            twod_images_report_folder (str): folder path to save twod data images.
            threed_images_report_folder (str): folder path to save threed data images.
        """

    def __init__(self,
                slide,
                general_input,
                metadb_3d_input,
                twod_images_report_folder,
                threed_images_report_folder) -> None:
        self.shapes = slide.shapes
        self.intrusion_areas = {"ROW 1":{},"ROW 2":{}}
        self.general_input = general_input
        self.metadb_3d_input = metadb_3d_input
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.logger = logging.getLogger("side_crash_logger")

    def intrusion_curve_format(self,source_window,curve,temporary_window_name,curve_name,target_curve = None):
        """
        This method is used to format intrusion curves.

        Args:
            source_window (str): source window name where the curve exists.
            curve (object): intrusion curve object.
            temporary_window_name (str): temporary window name.
            curve_name (str): custom curve name.
            target_curve(object):intrusion target Curve object


        Returns:
            int: 0 Always for Sucess.1 for Failure.
        """
        try:
            # self.logger.info("Intrusion curve fomat : {}".format(curve_name))
            # self.logger.info("")
            # starttime = datetime.now()
            # self.logger.info("Moving the {} curve to a Temporary window for custom formatting of title,yaxis,xaxis options and attibutes".format(curve_name))
            # self.logger.info(".")
            # self.logger.info(".")
            # self.logger.info(".")

            #moving the curve to a Temporary window
            if target_curve:
                utils.MetaCommand('xyplot curve copy "{}" {},{}'.format(source_window,curve.id,target_curve.id))
            else:
                utils.MetaCommand('xyplot curve copy "{}" {}'.format(source_window,curve.id))
            utils.MetaCommand('xyplot create "{}"'.format(temporary_window_name))
            if target_curve:
                utils.MetaCommand('xyplot curve paste "{}" 0 {},{}'.format(temporary_window_name,curve.id,target_curve.id))
            else:
                utils.MetaCommand('xyplot curve paste "{}" 0 {}'.format(temporary_window_name,curve.id))
            #getting temporary window object and its curves
            win = windows.Window(temporary_window_name, page_id=0)
            win.maximize()
            curve = win.get_curves('byname',name = curve.name)[0]
            #building x,y values for the curve to populate tables
            y_values = []
            for x in [0.03,0.04,0.05,0.06,0.07,0.08]:
                y_values.append(round(curve.get_y_values_from_x(specifier = 'first', xvalue =x)[0]))
            self.intrusion_areas[curve_name.rsplit(" ",1)[0]][curve_name.rsplit(" ",1)[1]] = y_values
            #applying custom style and size for plot title,xaxis,yaxis options
            utils.MetaCommand('xyplot gridoptions line major style "{}" 0 2'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions yaxis active "{}" 0 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions labels yposition "{}" 0 left'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions labels yalign "{}" 0 left'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions axyrange "{}" 0 0 0 600'.format(temporary_window_name))
            utils.MetaCommand('xyplot gridoptions yspace "{}" 0 40'.format(temporary_window_name))
            utils.MetaCommand('xyplot plotoptions title set "{}" 0 "{}"'.format(temporary_window_name,curve_name))
            utils.MetaCommand('xyplot axisoptions ylabel set "{}" 0 "Intrusion [mm]"'.format(temporary_window_name))
            utils.MetaCommand('xyplot curve select "{}" {}'.format(temporary_window_name,curve.id))
            if target_curve:
                utils.MetaCommand('xyplot curve deselect "{}" {}'.format(temporary_window_name,target_curve.id))
            utils.MetaCommand('xyplot curve set style "{}" selected 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot curve set linewidth "{}" selected 9'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions ylabel font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions labels yfont "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" 0 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions xaxis active "{}" 0 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions xlabel font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions labels xfont "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot plotoptions title font "{}" 0 "Arial,44,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions xaxis deactive "{}" 0 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot curve select "{}" all'.format(temporary_window_name))
            #endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into executive report slide:\n{}".format(e))
            self.logger.info("")
            return 1
        # self.logger.info("Intrusion curve format complete")
        # self.logger.info("Time Taken : {}".format(endtime - starttime))
        # self.logger.info("")
        return 0

    def edit(self):
        """
        This method is used to iterate all the shapes of the executive report slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """
        from pptx.util import Pt
        from PIL import Image

        try:
            self.logger.info("Started seeding data into executive report slide")
            self.logger.info("")
            starttime = datetime.now()
            #maximizing the survival space window
            if not self.general_input.survival_space_window_name in ["null","none",""]:
                #maximizing the survival space window
                survival_space_window_name = self.general_input.survival_space_window_name
                survival_space_window_obj = windows.WindowByName(survival_space_window_name)
                if survival_space_window_obj:
                    utils.MetaCommand('window maximize "{}"'.format(survival_space_window_name))
                    #sorting shapes of the slide based on their names
                    shapes = [shape for shape in self.shapes]
                    shapes.sort(key = lambda x:x.name)
                    #iterating through the shapes of the executive report slide
                    for shape in shapes:
                        #image insertion for the shape named "Image"
                        if shape.name == "Image":
                            initial_time_curve_name = "SS_0MS"
                            plot_id = 0
                            page_id=0
                            plot = plot2d.Plot(plot_id, survival_space_window_name, page_id)
                            title = plot.get_title()
                            org_name = title.get_text()
                            title.set_text("{} 0MS AND 150MS".format(org_name))
                            plot.activate()
                            #getting "Survival Space" plot object to activate and showing initial and final time curves
                            if not self.general_input.survival_space_final_time in ["null","none",""]:
                                final_time_curve_name = "SS_{}MS".format(round(float(self.general_input.survival_space_final_time)))
                                final_curves = plot2d.CurvesByName(survival_space_window_name, final_time_curve_name, 0)
                                if final_curves:
                                    final_curves[0].show()
                                else:
                                    self.logger.info("WARNING : Survival space window does not contain '{}' curve from META 2D variable {}. Please update.".format(final_time_curve_name,GeneralVarInfo.survival_space_final_time_key))
                                    self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.survival_space_final_time_key))
                                self.logger.info("")
                            initial_curve = plot2d.CurvesByName(survival_space_window_name, initial_time_curve_name, 0)[0]
                            initial_curve.show()
                            target_curves = plot2d.CurvesByName(survival_space_window_name, "*TARGET", 0)
                            _ret = [curve.show() if curve.id != 3 else curve.hide() for curve in target_curves]
                            #custom formatting for the "Survival Space" plot title,yaxis,xaxis options
                            utils.MetaCommand('xyplot axisoptions xaxis active "{}" {} 0'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot gridoptions xspace "{}" {} 20'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions axxrange "{}" {} 0 175 400'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot gridoptions line major style "{}" {} 0'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions xaxis active "{}" {} 0'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions xlabel font "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions labels xfont "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions xaxis deactive "{}" {} 0'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions yaxis active "{}" {} 0'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions yauto on "{}" {}'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions ylabel font "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions labels yfont "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" {} 0'.format(survival_space_window_name, plot_id))
                            utils.MetaCommand('xyplot plotoptions title font "{}" {} "Arial,10,-1,5,75,0,0,0,0,0"'.format(survival_space_window_name, plot_id))
                            #capturing "Survival Space" plot image
                            if self.twod_images_report_folder is not None:
                                image_path = os.path.join(self.twod_images_report_folder,survival_space_window_name+"_"+title.get_text()+".png").replace(" ","_")
                                capture_image(image_path,survival_space_window_name,shape.width,shape.height,transparent = True)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {},{} | SOURCE PLOT : {} | SOURCE WINDOW : {}".format(final_time_curve_name,initial_time_curve_name,title.get_text(),survival_space_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                transparent_image_path = image_path.replace(".png","")+"_transparent.png"
                                picture = self.shapes.add_picture(transparent_image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                title.set_text(org_name)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #removing transparent image
                                os.remove(transparent_image_path)
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #table population for the shape named "Table 1"
                        elif shape.name == "Table 1":
                            #getting row 2 object and inserting survival space final value in cell 0
                            if not self.general_input.survival_space_final_value not in ["null","none",""]:
                                final_time = round(float(self.general_input.survival_space_final_value))
                                text_frame_1 = shape.table.rows[2].cells[0].text_frame
                                font = text_frame_1.paragraphs[0].font
                                font.size = Pt(11)
                                text_frame_1.paragraphs[0].text = " {}".format(final_time)
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.survival_space_final_key))
                                self.logger.info("")
                            #getting row 2 object and inserting survival space peak value in cell 1
                            if not self.general_input.survival_space_peak_value not in ["null","none",""]:
                                text_frame_2 = shape.table.rows[2].cells[1].text_frame
                                font = text_frame_2.paragraphs[0].font
                                font.size = Pt(11)
                                text_frame_2.paragraphs[0].text = str(round(float(self.general_input.survival_space_peak_value)))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.survival_space_peak_key))
                                self.logger.info("")
                            #getting row 3 object and inserting peak time in cell 1
                            if self.general_input.peak_time_display_value not in ["null","none",""]:
                                text_frame_3 = shape.table.rows[3].cells[1].text_frame
                                font = text_frame_3.paragraphs[0].font
                                font.size = Pt(11)
                                text_frame_3.paragraphs[0].text = str(round(float(self.general_input.peak_time_display_value)))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.peak_time_display_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 3"
                        elif shape.name == "Image 3":
                            #visualizing "f28_front_door" critical part set to capture image at peak state
                            data = self.metadb_3d_input.critical_sections["f28_front_door"]
                            visualize_3d_critical_section(data)
                            utils.MetaCommand('grstyle scalarfringe enable')
                            utils.MetaCommand('0:options state variable "serial=1"')
                            utils.MetaCommand('options fringebar off')
                            if self.threed_images_report_folder:
                                image_path = os.path.join(self.threed_images_report_folder,"{}_{}.png".format(self.general_input.threed_window_name,data["name"])).replace(" ","_")
                                capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height)
                                self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                                self.logger.info("SOURCE MODEL : 0")
                                self.logger.info("STATE : PEAK STATE")
                                self.logger.info("PID NAME SHOW FILTER : {} ".format(data["hes"] if "hes" in data.keys() else "null"))
                                self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                                self.logger.info("PID NAME ERASE FILTER : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                                self.logger.info("PID'S TO ERASE : {} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null"))
                                self.logger.info("ERASE BOX : {} ".format(data["erase_box"] if "erase_box" in data.keys() else "null"))
                                self.logger.info("IMAGE VIEW : {} ".format(data["view"] if "view" in data.keys() else "null"))
                                self.logger.info("TRANSPARENCY LEVEL : 50" )
                                self.logger.info("TRANSPARENT PID'S : {} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null"))
                                self.logger.info("COMP NAME : {} ".format(data["name"] if "name" in data.keys() else "null"))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT MODEL IMAGES :")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                utils.MetaCommand('color pid transparency reset act')
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 4"
                        elif shape.name == "Image 4":
                            #visualizing "f28_rear_door" critical part set to capture image at peak state
                            data = self.metadb_3d_input.critical_sections["f28_rear_door"]
                            visualize_3d_critical_section(data)
                            utils.MetaCommand('grstyle scalarfringe enable')
                            utils.MetaCommand('0:options state variable "serial=1"')
                            utils.MetaCommand('options fringebar off')
                            if self.threed_images_report_folder:
                                image_path = os.path.join(self.threed_images_report_folder,"{}_{}.png".format(self.general_input.threed_window_name,data["name"])).replace(" ","_")
                                capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height)
                                self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                                self.logger.info("SOURCE MODEL : 0")
                                self.logger.info("STATE : PEAK STATE")
                                self.logger.info("PID NAME SHOW FILTER : {} ".format(data["hes"] if "hes" in data.keys() else "null"))
                                self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                                self.logger.info("PID NAME ERASE FILTER : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                                self.logger.info("PID'S TO ERASE : {} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null"))
                                self.logger.info("ERASE BOX : {} ".format(data["erase_box"] if "erase_box" in data.keys() else "null"))
                                self.logger.info("IMAGE VIEW : {} ".format(data["view"] if "view" in data.keys() else "null"))
                                self.logger.info("TRANSPARENCY LEVEL : 50" )
                                self.logger.info("TRANSPARENT PID'S : {} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null"))
                                self.logger.info("COMP NAME : {} ".format(data["name"] if "name" in data.keys() else "null"))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT MODEL IMAGES :")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                utils.MetaCommand('color pid transparency reset act')
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 2"
                        elif shape.name == "Image 2":
                            #maximizing the metapost window
                            utils.MetaCommand('window maximize "{}"'.format(self.general_input.threed_window_name))
                            #visualizing "f21_upb_inner" critical part set
                            data = self.metadb_3d_input.critical_sections["f21_upb_inner"]
                            visualize_3d_critical_section(data)
                            #utils.MetaCommand('color pid transparency reset act')
                            utils.MetaCommand('grstyle scalarfringe enable')
                            #adding a yz plane and slicing the critiical part set with width 500
                            utils.MetaCommand('plane new DEFAULT_PLANE_YZ xyz 1657.996826,-16.504395,576.072754 1,0,0')
                            utils.MetaCommand('plane edit perpendicular 0/1/0 DEFAULT_PLANE_YZ')
                            utils.MetaCommand('plane toggleopts enable sectionclip DEFAULT_PLANE_YZ')
                            utils.MetaCommand('plane toggleopts disable stateauto DEFAULT_PLANE_YZ')
                            utils.MetaCommand('plane options cut autovisible DEFAULT_PLANE_YZ')
                            utils.MetaCommand('plane options onlysection enable DEFAULT_PLANE_YZ')
                            utils.MetaCommand('plane toggleopts enable clip DEFAULT_PLANE_YZ')
                            utils.MetaCommand('plane toggleopts disable sectionclip DEFAULT_PLANE_YZ')
                            utils.MetaCommand('plane toggleopts enable slice DEFAULT_PLANE_YZ')
                            utils.MetaCommand('plane options slicewidth 500.000000 DEFAULT_PLANE_YZ')
                            utils.MetaCommand('grstyle scalarfringe enable')
                            utils.MetaCommand('view default front')
                            utils.MetaCommand('0:options state variable "serial=1"')
                            utils.MetaCommand('options fringebar off')
                            #capturing "f21_upb_inner" image at peak state
                            if self.threed_images_report_folder:
                                image_path = os.path.join(self.threed_images_report_folder,"{}_{}.png".format(self.general_input.threed_window_name,"F21_UPBPILLAR_AT_PEAK_STATE")).replace(" ","_")
                                capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height,view = "front")
                                self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                                self.logger.info("SOURCE MODEL : 0")
                                self.logger.info("STATE : PEAK STATE")
                                self.logger.info("PID NAME SHOW FILTER : {} ".format(data["hes"] if "hes" in data.keys() else "null"))
                                self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                                self.logger.info("PID NAME ERASE FILTER : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                                self.logger.info("PID'S TO ERASE : {} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null"))
                                self.logger.info("ERASE BOX : {} ".format(data["erase_box"] if "erase_box" in data.keys() else "null"))
                                self.logger.info("IMAGE VIEW : {} ".format(data["view"] if "view" in data.keys() else "null"))
                                self.logger.info("TRANSPARENCY LEVEL : 50" )
                                self.logger.info("TRANSPARENT PID'S : {} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null"))
                                self.logger.info("PLANE CUT & SLICE WIDTH: DEFAULT_PLANE_YZ & 500" )
                                self.logger.info("COMP NAME : {} ".format(data["name"] if "name" in data.keys() else "null"))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT MODEL IMAGES :")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting created plane and reverting transarency on critical part set visualization
                                utils.MetaCommand('plane delete DEFAULT_PLANE_YZ')
                                utils.MetaCommand('color pid transparency reset act')
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 5"
                        elif shape.name == "Image 5":
                            #setting temporary window name and maximizing front door accel window
                            temporary_window_name = "Temporary"
                            front_door_accel_window_name = self.general_input.front_door_accel_window_name
                            utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                            if self.general_input.front_shoulder_intrusion_curve_name not in ["null","none",""]:
                                front_shoulder_intrusion_curve_name = self.general_input.front_shoulder_intrusion_curve_name
                                #getting front shoulder intrusion curve object and modifying it to capture image
                                curves = plot2d.CurvesByName(front_door_accel_window_name, front_shoulder_intrusion_curve_name, 1)
                                if curves:
                                    curves[0].show()
                                    self.intrusion_curve_format(front_door_accel_window_name,curves[0],temporary_window_name,"ROW 1 SHOULDER")
                                else:
                                    self.logger.info("WARNING : Front Door Accel window does not contain '{}' curve from META 2D variable {}. Please update.".format(front_shoulder_intrusion_curve_name,GeneralVarInfo.front_shoulder_intrusion_curve_key))
                                    self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.front_shoulder_intrusion_curve_key))
                                self.logger.info("")
                            #capturing front shoulder intrusion curve plot image
                            if self.twod_images_report_folder:
                                image_path = os.path.join(self.twod_images_report_folder,front_door_accel_window_name+"_"+"ROW 1 SHOULDER"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(front_shoulder_intrusion_curve_name,self.general_input.front_shoulder_intrusion_curve_key,front_door_accel_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting temporary window
                                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 6"
                        elif shape.name == "Image 6":
                            #setting temporary window name and maximizing front door accel window
                            temporary_window_name = "Temporary"
                            front_door_accel_window_name = self.general_input.front_door_accel_window_name
                            utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                            if self.general_input.front_abdomen_intrusion_curve_name not in ["null","none",""]:
                                front_abdomen_intrusion_curve_name = self.general_input.front_abdomen_intrusion_curve_name
                                #getting front abdomen intrusion curve object and modifying it to capture image
                                curves = plot2d.CurvesByName(front_door_accel_window_name, front_abdomen_intrusion_curve_name, 1)
                                if curves:
                                    curves[0].show()
                                    self.intrusion_curve_format(front_door_accel_window_name,curves[0],temporary_window_name,"ROW 1 ABDOMEN")
                                else:
                                    self.logger.info("WARNING : Front Door Accel window does not contain '{}' curve from META 2D variable {}. Please update.".format(front_abdomen_intrusion_curve_name,GeneralVarInfo.front_abdomen_intrusion_curve_key))
                                    self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.front_abdomen_intrusion_curve_key))
                                self.logger.info("")
                            #capturing front abdomen intrusion curve plot image
                            if self.twod_images_report_folder:
                                image_path = os.path.join(self.twod_images_report_folder,front_door_accel_window_name+"_"+"ROW 1 ABDOMEN"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(front_abdomen_intrusion_curve_name,self.general_input.front_abdomen_intrusion_curve_key,front_door_accel_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting temporary window
                                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 8"
                        elif shape.name == "Image 8":
                            #setting temporary window name and maximizing front door accel window
                            temporary_window_name = "Temporary"
                            front_door_accel_window_name = self.general_input.front_door_accel_window_name
                            utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                            if self.general_input.front_femur_intrusion_curve_name not in ["null","none",""]:
                                front_femur_intrusion_curve_name = self.general_input.front_femur_intrusion_curve_name
                                #getting front femur intrusion curve object and modifying it to capture image
                                curves = plot2d.CurvesByName(front_door_accel_window_name, front_femur_intrusion_curve_name, 1)
                                if curves:
                                    curves[0].show()
                                    self.intrusion_curve_format(front_door_accel_window_name,curves[0],temporary_window_name,"ROW 1 FEMUR")
                                else:
                                    self.logger.info("WARNING : Front Door Accel window does not contain '{}' curve from META 2D variable {}. Please update.".format(front_femur_intrusion_curve_name,GeneralVarInfo.front_femur_intrusion_curve_key))
                                    self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.front_femur_intrusion_curve_key))
                                self.logger.info("")
                            #capturing front femur intrusion curve plot image
                            if self.twod_images_report_folder:
                                image_path = os.path.join(self.twod_images_report_folder,front_door_accel_window_name+"_"+"ROW 1 FEMUR"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(front_femur_intrusion_curve_name,self.general_input.front_femur_intrusion_curve_key,front_door_accel_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting temporary window
                                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 7"
                        elif shape.name == "Image 7":
                            #setting temporary window name and maximizing front door accel window
                            temporary_window_name = "Temporary"
                            front_door_accel_window_name = self.general_input.front_door_accel_window_name
                            utils.MetaCommand('window maximize {}'.format(front_door_accel_window_name))
                            if self.general_input.front_pelvis_intrusion_curve_name not in ["null","none",""]:
                                front_pelvis_intrusion_curve_name = self.general_input.front_pelvis_intrusion_curve_name
                                #getting front pelvis intrusion curve object and modifying it to capture image
                                curves = plot2d.CurvesByName(front_door_accel_window_name, front_pelvis_intrusion_curve_name, 1)
                                if curves:
                                    curves[0].show()
                                    self.intrusion_curve_format(front_door_accel_window_name,curves[0],temporary_window_name,"ROW 1 PELVIS")
                                else:
                                    self.logger.info("WARNING : Front Door Accel window does not contain '{}' curve from META 2D variable {}. Please update.".format(front_pelvis_intrusion_curve_name,GeneralVarInfo.front_pelvis_intrusion_curve_key))
                                    self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.front_pelvis_intrusion_curve_key))
                                self.logger.info("")
                            #capturing front pelvis intrusion curve plot image
                            if self.twod_images_report_folder:
                                image_path = os.path.join(self.twod_images_report_folder,front_door_accel_window_name+"_"+"ROW 1 PELVIS"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(front_pelvis_intrusion_curve_name,self.general_input.front_pelvis_intrusion_curve_key,front_door_accel_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting temporary window
                                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 0"
                        elif shape.name == "Image 0":
                            #setting temporary window name and maximizing rear door accel window
                            temporary_window_name = "Temporary"
                            rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                            utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                            if self.general_input.rear_shoulder_intrusion_curve_name not in ["null","none",""]:
                                rear_shoulder_intrusion_curve_name = self.general_input.rear_shoulder_intrusion_curve_name
                                #getting rear shoulder intrusion curve object and modifying it to capture image
                                curves = plot2d.CurvesByName(rear_door_accel_window_name, rear_shoulder_intrusion_curve_name, 1)
                                if curves:
                                    curves[0].show()
                                    target_curve = plot2d.CurvesByName(rear_door_accel_window_name, "*TARGET", 0)[0]
                                    target_curve.show()
                                    self.intrusion_curve_format(rear_door_accel_window_name,curves[0],temporary_window_name,"ROW 2 SHOULDER",target_curve=target_curve)
                                else:
                                    self.logger.info("WARNING : Rear Door Accel window does not contain '{}' curve from META 2D variable {}. Please update.".format(rear_shoulder_intrusion_curve_name,GeneralVarInfo.rear_shoulder_intrusion_curve_key))
                                    self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.rear_shoulder_intrusion_curve_key))
                                self.logger.info("")
                            #capturing rear shoulder intrusion curve plot image
                            if self.twod_images_report_folder:
                                image_path = os.path.join(self.twod_images_report_folder,rear_door_accel_window_name+"_"+"ROW 2 SHOULDER"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(rear_shoulder_intrusion_curve_name,self.general_input.rear_shoulder_intrusion_curve_key,rear_door_accel_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting temporary window
                                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 10"
                        elif shape.name == "Image 10":
                            #setting temporary window name and maximizing rear door accel window
                            temporary_window_name = "Temporary"
                            rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                            utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                            if self.general_input.rear_abdomen_intrusion_curve_name not in ["null","none",""]:
                                rear_abdomen_intrusion_curve_name = self.general_input.rear_abdomen_intrusion_curve_name
                                #getting rear abdomen intrusion curve object and modifying it to capture image
                                curves = plot2d.CurvesByName(rear_door_accel_window_name, rear_abdomen_intrusion_curve_name, 1)
                                if curves:
                                    curves[0].show()
                                    target_curve = plot2d.CurvesByName(rear_door_accel_window_name, "*TARGET", 0)[0]
                                    target_curve.show()
                                    self.intrusion_curve_format(rear_door_accel_window_name,curves[0],temporary_window_name,"ROW 2 ABDOMEN",target_curve=target_curve)
                                else:
                                    self.logger.info("WARNING : Rear Door Accel window does not contain '{}' curve from META 2D variable {}. Please update.".format(rear_abdomen_intrusion_curve_name,GeneralVarInfo.rear_abdomen_intrusion_curve_key))
                                    self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.rear_abdomen_intrusion_curve_key))
                                self.logger.info("")
                            #capturing rear abdomen intrusion curve plot image
                            if self.twod_images_report_folder:
                                image_path = os.path.join(self.twod_images_report_folder,rear_door_accel_window_name+"_"+"ROW 2 ABDOMEN"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(rear_abdomen_intrusion_curve_name,self.general_input.rear_abdomen_intrusion_curve_key,rear_door_accel_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting temporary window
                                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 12"
                        elif shape.name == "Image 12":
                            #setting temporary window name and maximizing rear door accel window
                            temporary_window_name = "Temporary"
                            rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                            utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                            if self.general_input.rear_femur_intrusion_curve_name not in ["null","none",""]:
                                rear_femur_intrusion_curve_name = self.general_input.rear_femur_intrusion_curve_name
                                #getting rear femur intrusion curve object and modifying it to capture image
                                curves = plot2d.CurvesByName(rear_door_accel_window_name, rear_femur_intrusion_curve_name, 1)
                                if curves:
                                    curves[0].show()
                                    target_curve = plot2d.CurvesByName(rear_door_accel_window_name, "*TARGET", 0)[0]
                                    target_curve.show()
                                    self.intrusion_curve_format(rear_door_accel_window_name,curves[0],temporary_window_name,"ROW 2 FEMUR",target_curve=target_curve)
                                else:
                                    self.logger.info("WARNING : Rear Door Accel window does not contain '{}' curve from META 2D variable {}. Please update.".format(rear_femur_intrusion_curve_name,GeneralVarInfo.rear_femur_intrusion_curve_key))
                                    self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.rear_femur_intrusion_curve_key))
                                self.logger.info("")
                            #capturing rear femur intrusion curve plot image
                            if self.twod_images_report_folder:
                                image_path = os.path.join(self.twod_images_report_folder,rear_door_accel_window_name+"_"+"ROW 2 FEMUR"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(rear_femur_intrusion_curve_name,self.general_input.rear_femur_intrusion_curve_key,rear_door_accel_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting temporary window
                                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                        #image insertion for the shape named "Image 11"
                        elif shape.name == "Image 11":
                            #setting temporary window name and maximizing rear door accel window
                            temporary_window_name = "Temporary"
                            rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                            utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                            if self.general_input.rear_pelvis_intrusion_curve_name not in ["null","none",""]:
                                rear_pelvis_intrusion_curve_name = self.general_input.rear_pelvis_intrusion_curve_name
                                #getting rear pelvis intrusion curve object and modifying it to capture image
                                curves = plot2d.CurvesByName(rear_door_accel_window_name, rear_pelvis_intrusion_curve_name, 1)
                                if curves:
                                    curves[0].show()
                                    target_curve = plot2d.CurvesByName(rear_door_accel_window_name, "*TARGET", 0)[0]
                                    target_curve.show()
                                    self.intrusion_curve_format(rear_door_accel_window_name,curves[0],temporary_window_name,"ROW 2 PELVIS",target_curve=target_curve)
                                else:
                                        self.logger.info("WARNING : Rear Door Accel window does not contain '{}' curve from META 2D variable {}. Please update.".format(rear_femur_intrusion_curve_name,GeneralVarInfo.rear_femur_intrusion_curve_key))
                                        self.logger.info("")
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.rear_femur_intrusion_curve_key))
                                self.logger.info("")
                            #capturing rear pelvis intrusion curve plot image
                            if self.twod_images_report_folder:
                                image_path = os.path.join(self.twod_images_report_folder,rear_door_accel_window_name+"_"+"ROW 2 PELVIS"+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(rear_pelvis_intrusion_curve_name,self.general_input.rear_pelvis_intrusion_curve_key,rear_door_accel_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #deleting temporary window
                                utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                    #iterating through table shapes
                    for shape in self.shapes:
                        #table population for the shape named "Table 2"
                        if shape.name == "Table 2":
                            #getting the table object and setting some temporary variables to iterate through the cells of the table automatically
                            table = shape.table
                            row_index = 2
                            node_index = 0
                            #adding a new row to the tables
                            add_row(table)
                            #iterating through the x,y values of each ROW 1 nodeS built from intrusion curve format method
                            for iindex,(key,values) in enumerate(self.intrusion_areas["ROW 1"].items()):
                                #adding a new row for the third node x,y values
                                if iindex == 2:
                                    add_row(table)
                                    row_index = row_index+1
                                #getting index for all four node x,y values
                                if iindex == 1 or iindex == 3:
                                    node_index = 7
                                else:
                                    node_index = 0
                                #inserting node name into the index cell of the node
                                text_frame_1 = table.rows[row_index].cells[node_index].text_frame
                                font = text_frame_1.paragraphs[0].font
                                font.bold = True
                                font.size = Pt(11)
                                text_frame_1.paragraphs[0].text = key.capitalize()
                                #inserting  x,y values of the nodes
                                for index,value in enumerate(values):
                                    text_frame = table.rows[row_index].cells[node_index+index+1].text_frame
                                    font = text_frame.paragraphs[0].font
                                    font.size = Pt(11)
                                    text_frame.paragraphs[0].text = str(value)
                        #table population for the shape named "Table 2"
                        elif shape.name == "Table 3":
                            #getting the table object and setting some temporary variables to iterate through the cells of the table automatically
                            table = shape.table
                            row_index = 2
                            node_index = 0
                            #adding a new row to the tables
                            add_row(table)
                            #iterating through the x,y values of each ROW 2 nodes built from intrusion curve format method
                            for iindex,(key,values) in enumerate(self.intrusion_areas["ROW 2"].items()):
                                #adding a new row for the third node x,y values
                                if iindex == 2:
                                    add_row(table)
                                    row_index = row_index+1
                                #getting index for all four node x,y values
                                if iindex == 1 or iindex == 3:
                                    node_index = 7
                                else:
                                    node_index = 0
                                #inserting node name into the index cell of the node
                                text_frame_1 = table.rows[row_index].cells[node_index].text_frame
                                font = text_frame_1.paragraphs[0].font
                                font.bold = True
                                font.size = Pt(11)
                                text_frame_1.paragraphs[0].text = key.capitalize()
                                #inserting  x,y values of the nodes
                                for index,value in enumerate(values):
                                    text_frame = table.rows[row_index].cells[node_index+index+1].text_frame
                                    font = text_frame.paragraphs[0].font
                                    font.size = Pt(11)
                                    text_frame.paragraphs[0].text = str(value)
                else:
                    self.logger.info("ERROR : 2D METADB does not contain 'Survival Space' window. Please update.")
            else:
                self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.survival_space_window_key))
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into executive report slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into executive report slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")
        return 0
