# PYTHON script
"""
This script is used for all the automation process of CBU and Barrier position slide of thesis report.
"""

import os
import logging
from datetime import datetime

from meta import utils
from meta import nodes
from meta import models
from meta import plot2d

from src.meta_utilities import capture_image
from src.meta_utilities import visualize_3d_critical_section
from src.metadb_info import GeneralVarInfo

class CBUAndBarrierPositionSlide():
    """
        This class is used to automate the CBU and Barrier position slide of thesis report.

        Args:
            slide (object): cae quality pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
            metadb_3d_input (Meta3DInfo): Meta3DInfo class object.
            threed_images_report_folder (str): folder path to save threed data images.
        """
    def __init__(self,
                slide,
                general_input,
                metadb_3d_input,
                threed_images_report_folder) -> None:
        self.shapes = slide.shapes
        self.general_input = general_input
        self.metadb_3d_input = metadb_3d_input
        self.threed_images_report_folder = threed_images_report_folder
        self.logger = logging.getLogger("side_crash_logger")

    def edit(self):
        """
        This method is used to iterate all the shapes of the cbu and barrier position slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """

        from PIL import Image
        from pptx.util import Pt

        try:
            self.logger.info("Started seeding data into cbu and barrier position slide")
            self.logger.info("")
            starttime = datetime.now()
            #maximizing metapost window
            utils.MetaCommand('window maximize "{}"'.format(self.general_input.threed_window_name))
            utils.MetaCommand('0:options state original')
            utils.MetaCommand('options fringebar off')
            #iterating through the shapes of the cbu and barrier position slide
            for shape in self.shapes:
                #image insertion for the shape named "Image 4"
                if shape.name == "Image 4":
                    #visualizing all the critical part sets to capture whole cbu and barrier image at original state
                    critical_data = self.metadb_3d_input.critical_sections
                    for (index,(_critical_section,value)) in enumerate(critical_data.items()):
                        and_filter = False
                        if index>0:
                            and_filter = True
                        visualize_3d_critical_section(value,and_filter = and_filter)
                    # If twod_images_report_folder is not None the capture image and add the picture into slide
                    if self.threed_images_report_folder is not None:
                        image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"CBU_AND_BARRIER_TOP_VIEW"+".png").replace(" ","_")
                        capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height,rotate = Image.ROTATE_90,view = "top",transparent=True)
                        self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                        self.logger.info("")
                        self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                        self.logger.info("SOURCE MODEL : 0")
                        self.logger.info("STATE : ORIGINAL STATE")
                        self.logger.info("PID NAME SHOW FILTER : {} ".format("CBU"))
                        self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format("CBU"))
                        self.logger.info("IMAGE VIEW : TOP ")
                        self.logger.info("COMP NAME : CBU AND BARRIER ")
                        self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                        self.logger.info("OUTPUT MODEL IMAGES :")
                        self.logger.info(image_path)
                        self.logger.info("")
                        #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                        transparent_image_path = image_path.replace(".png","")+"_transparent.png"
                        picture = self.shapes.add_picture(transparent_image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                        picture.crop_left = 0
                        picture.crop_right = 0
                        #removing transparent image
                        os.remove(transparent_image_path)
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                        self.logger.info("")
                #image insertion for the shape named "Image 3"
                elif shape.name == "Image 3":
                   #visualizing all the critical part sets to capture whole cbu and barrier image at original state with gray color
                    critical_data = self.metadb_3d_input.critical_sections
                    for (index,(_critical_section,value)) in enumerate(critical_data.items()):
                        and_filter = False
                        if index>0:
                            and_filter = True
                        visualize_3d_critical_section(value,and_filter = and_filter)
                    utils.MetaCommand('color pid Gray act')
                    # If twod_images_report_folder is not None the capture image and add the picture into slide
                    if self.threed_images_report_folder is not None:
                        image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"CBU_AND_BARRIER_IMPACT_LEFT_VIEW"+".png").replace(" ","_")
                        capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height,view = "left",transparent=True)
                        self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                        self.logger.info("")
                        self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                        self.logger.info("SOURCE MODEL : 0")
                        self.logger.info("STATE : ORIGINAL STATE")
                        self.logger.info("PID NAME SHOW FILTER : {} ".format("CBU"))
                        self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format("CBU"))
                        self.logger.info("IMAGE VIEW : TOP ")
                        self.logger.info("COMP NAME : CBU AND BARRIER IMPACT ")
                        self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                        self.logger.info("OUTPUT MODEL IMAGES :")
                        self.logger.info(image_path)
                        self.logger.info("")
                        #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                        transparent_image_path = image_path.replace(".png","")+"_transparent.png"
                        picture = self.shapes.add_picture(transparent_image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                        picture.crop_left = 0
                        picture.crop_right = 0
                        #removing transparent image
                        os.remove(transparent_image_path)
                        #reverting color
                        utils.MetaCommand('color pid reset act')
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                        self.logger.info("")
                #table population for the shape named "Table 4"
                elif shape.name == "Table 4":
                    #getting the table object
                    table_obj = shape.table
                    # If test_mass_value not in null,none,""
                    if self.general_input.test_mass_value not in ["null","none",""]:
                        #getting row 1 object and inserting test mass value in cell 1
                        text_frame = table_obj.rows[1].cells[1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        font.bold = True
                        font.underline = True
                        text_frame.paragraphs[0].text = str(round(float(self.general_input.test_mass_value)*1000,2))
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.test_mass_key))
                        self.logger.info("")
                    if self.general_input.physical_mass_value not in ["null","none",""]:
                        #getting row 2 object and inserting physical mass value in cell 1
                        text_frame = table_obj.rows[2].cells[1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(round(float(self.general_input.physical_mass_value)*1000,2))
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.physical_mass_key))
                        self.logger.info("")
                    if self.general_input.added_mass_value not in ["null","none",""]:
                        #getting row 3 object and inserting added mass value in cell 1
                        text_frame = table_obj.rows[3].cells[1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(round(float(self.general_input.added_mass_value)*1000,2))
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.added_mass_key))
                        self.logger.info("")
                    if self.general_input.total_mass_value not in ["null","none",""]:
                        #getting row 6 object and inserting total mass value in cell 1
                        text_frame = table_obj.rows[6].cells[1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(round(float(self.general_input.total_mass_value)*1000,2))
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.total_mass_key))
                        self.logger.info("")

                    if self.general_input.test_mass_value not in ["null","none",""] and self.general_input.total_mass_value not in ["null","none",""]:
                        #getting row 7 object and inserting test mass - total mass value in cell 1
                        text_frame = table_obj.rows[7].cells[1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(round((float(self.general_input.test_mass_value)-float(self.general_input.total_mass_value))*1000,2))
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' '{}' is not available or invalid. Please update.".format(GeneralVarInfo.total_mass_key,GeneralVarInfo.test_mass_key))
                        self.logger.info("")
                #table population for the shape named "Table 1"
                elif shape.name == "Table 1":
                    #getting the table object
                    table_obj = shape.table
                    # If MDB_fr_node_id not in null,none,""
                    if self.general_input.MDB_fr_node_id not in ["null","none",""]:
                        #getting row 3 object and inserting MDB front node x value
                        MDB_fr_node_id = int(self.general_input.MDB_fr_node_id)
                        MDB_fr_node = nodes.Node(id=MDB_fr_node_id, model_id=0)
                        text_frame = table_obj.rows[3].cells[1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(round(MDB_fr_node.x))
                        #getting row 2 cell 1 text
                        text_frame = table_obj.rows[2].cells[1].text_frame
                        target_z_value = int(text_frame.paragraphs[0].text.strip())
                        #getting row 4 object and inserting MDB font node x - above z value in cell 1
                        text_frame = table_obj.rows[4].cells[1].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        value = round(MDB_fr_node.x) - target_z_value
                        text_frame.paragraphs[0].text = str("+"+str(value) if value>0 else value)
                    else:
                        self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.MDB_fr_node_key))
                        self.logger.info("")
                #table population for the shape named "Table 2"
                elif shape.name == "Table 2":
                    #getting current model resultset,suspension nodes,MDB front and rear node objects
                    model = models.Model(0)
                    res = model.get_current_resultset()
                    if all(var not in ["null","none",""] for var in [self.general_input.struck_subframe_node_ids,self.general_input.MDB_fr_node_id,self.general_input.MDB_rr_node_id]):
                        struck_subframe_node_ids = self.general_input.struck_subframe_node_ids
                        target_nodes = struck_subframe_node_ids.split("/")
                        curve_types = plot2d.CurvesTypesDynaWithNames(self.general_input.binout_directory)
                        id_sorted_search_nodes = []
                        search_target = []
                        identified = []
                        for one_type in curve_types:
                            types = one_type[0]
                            if (types == 'nodout-Node'):
                                entities = one_type[1]
                                for one_entity in entities:
                                    logged = 0
                                    entity_id = one_entity[0]
                                    entity_name = one_entity[1]
                                    if (entity_name in target_nodes):
                                        is_it_real = target_nodes.index(entity_name)
                                        id_sorted_search_nodes.append(str(entity_id))
                                        search_target.append(entity_name)
                                        identified.append(target_nodes[is_it_real])
                                        logged = 1
                                    elif str(entity_id) in target_nodes and (logged == 0):
                                        is_it_real = target_nodes.index(entity_id)
                                        id_sorted_search_nodes.append(target_nodes[is_it_real])
                                        search_target.append(str(entity_id))
                                        identified.append(target_nodes[is_it_real])
                        struck_subframe_node1 = nodes.Node(id=int(target_nodes[0]), model_id=0)
                        struck_subframe_node2 = nodes.Node(id=int(target_nodes[1]), model_id=0)
                        MDB_fr_node_id = int(self.general_input.MDB_fr_node_id)
                        MDB_fr_node = nodes.Node(id=MDB_fr_node_id, model_id=0)
                        MDB_rr_node_id = int(self.general_input.MDB_rr_node_id)
                        MDB_rr_node = nodes.Node(id=MDB_rr_node_id, model_id=0)
                        #calculating distance form suspension node 1 to MDB front and MDB rear node
                        distance_nfr_s1 = MDB_fr_node.get_distance_from_node(res, struck_subframe_node1, res)
                        distance_nrr_s1 = MDB_fr_node.get_distance_from_node(res, struck_subframe_node1, res)
                        distance_nfr_s1 = sum([dist**2 for dist in distance_nfr_s1])
                        distance_nrr_s1 = sum([dist**2 for dist in distance_nrr_s1])
                        #identifying suspension front and rear nodes based on the ditance calculated above
                        if distance_nfr_s1>distance_nrr_s1:
                            suspension_rr_node = struck_subframe_node1
                            suspension_fr_node = struck_subframe_node2
                        else:
                            suspension_rr_node = struck_subframe_node2
                            suspension_fr_node = struck_subframe_node1
                        #getting table object
                        table_obj = shape.table
                        #getting row 2 object and inserting MDB front node x - Suspension front node x value in cell 2
                        text_frame = table_obj.rows[2].cells[2].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(abs(round(MDB_fr_node.x) - round(suspension_fr_node.x)))
                        #getting row 2 cell 1 text
                        text_frame = table_obj.rows[2].cells[1].text_frame
                        front_target_overlap = int(text_frame.paragraphs[0].text.strip())
                        #getting row 2 object and inserting front target overlap difference in cell 3
                        text_frame = table_obj.rows[2].cells[3].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        value = abs(round(MDB_fr_node.x) - round(suspension_fr_node.x)) - front_target_overlap
                        text_frame.paragraphs[0].text = str("+"+str(value) if value>0 else value)
                        #getting row 3 object and inserting MDB rear node x - Suspension rear node x value in cell 2
                        text_frame = table_obj.rows[3].cells[2].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        text_frame.paragraphs[0].text = str(abs(round(MDB_rr_node.x) - round(suspension_rr_node.x)))
                        #getting row 3 cell 1 text
                        text_frame = table_obj.rows[3].cells[1].text_frame
                        rear_target_overlap = int(text_frame.paragraphs[0].text.strip())
                        #getting row 3 object and inserting rear target overlap difference in cell 3
                        text_frame = table_obj.rows[3].cells[3].text_frame
                        font = text_frame.paragraphs[0].font
                        font.name = 'Arial'
                        font.size = Pt(11)
                        value = abs(round(MDB_rr_node.x) - round(suspension_rr_node.x)) - rear_target_overlap
                        text_frame.paragraphs[0].text = str("+"+str(value) if value>0 else value)
                    else:
                        self.logger.info("ERROR : META 2D variables '{},{},{}' are not available or invalid. Please update.".format(GeneralVarInfo.struck_subframe_node_key,GeneralVarInfo.MDB_fr_node_key,GeneralVarInfo.MDB_rr_node_key))
                        self.logger.info("")
            utils.MetaCommand('options fringebar on')
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into cbu and barrier position slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into cbu and barrier position slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")

        return 0

    def revert(self):
        """
        revert _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """

        return 0
