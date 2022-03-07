# PYTHON script
"""
    _summary_

_extended_summary_

Returns:
    _type_: _description_
"""

import os

from meta import utils,parts,constants

from src.meta_utilities import capture_resized_image,visualize_3d_critical_section
from src.general_utilities import add_row

class BOMF21ROOFSlide():
    def __init__(self,
                slide,
                windows,
                general_input,
                metadb_3d_input,
                template_file,
                threed_images_report_folder,
                ppt_report_folder) -> None:
        self.shapes = slide.shapes
        self.windows = windows
        self.general_input = general_input
        self.metadb_3d_input = metadb_3d_input
        self.template_file = template_file
        self.threed_images_report_folder = threed_images_report_folder
        self.ppt_report_folder = ppt_report_folder
    def setup(self):
        """
        setup _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """

        return 0
    def edit(self, ):
        from PIL import Image
        from pptx.util import Pt
        self.setup()
        utils.MetaCommand('0:options state original')
        for shape in self.shapes:
            if shape.name == "Image 1":
                data = self.metadb_3d_input.critical_sections["f21_roof"]
                visualize_3d_critical_section(data)
                utils.MetaCommand('window maximize "MetaPost"')
                utils.MetaCommand('options fringebar off')
                image_path = os.path.join(self.threed_images_report_folder,"MetaPost"+"_"+"f21_roof".lower()+".png")
                capture_resized_image("MetaPost",shape.width,shape.height,image_path,rotate = Image.ROTATE_270,view = "btm")
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Table 1":
                data = self.metadb_3d_input.critical_sections["f21_roof"]
                prop_names = data["hes"]
                re_props = prop_names.split(",")
                entities = []
                for re_prop in re_props:
                    utils.MetaCommand('window maximize "MetaPost"')
                    entities.extend(self.metadb_3d_input.get_props(re_prop))
                table_obj = shape.table
                for id,prop in enumerate(entities):
                    part_type = parts.StringPartType(prop.type)
                    if part_type == "PSHELL":
                        part = parts.Part(id=prop.id,type = constants.PSHELL, model_id=0)
                        materials = part.get_materials('all')
                    elif part_type == "PSOLID":
                        part = parts.Part(id=prop.id,type = constants.PSOLID, model_id=0)
                        materials = part.get_materials('all')

                    add_row(table_obj)
                    prop_row = table_obj.rows[id+1]
                    text_frame = prop_row.cells[0].text_frame
                    font = text_frame.paragraphs[0].font
                    font.size = Pt(8)
                    text_frame.paragraphs[0].text = str(prop.id)

                    text_frame_name = prop_row.cells[1].text_frame
                    font_name = text_frame_name.paragraphs[0].font
                    font_name.size = Pt(8)
                    text_frame_name.paragraphs[0].text = str(prop.name)

                    text_frame_material = prop_row.cells[2].text_frame
                    font_material = text_frame_material.paragraphs[0].font
                    font_material.size = Pt(8)

                    text_frame_material.paragraphs[0].text = str(materials[0].name)

                    text_frame_thickness = prop_row.cells[3].text_frame
                    font_thickness = text_frame_thickness.paragraphs[0].font
                    font_thickness.size = Pt(8)
                    thickness = round(prop.shell_thick,1)
                    text_frame_thickness.paragraphs[0].text = str(thickness)
        self.revert()
        return 0
    def revert(self):
        """
        revert _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """
        utils.MetaCommand('0:options state variable "serial=1"')

        return 0
