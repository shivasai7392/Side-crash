# PYTHON script
"""
    _summary_

_extended_summary_

Returns:
    _type_: _description_
"""

import os

from meta import utils
from meta import models
from meta import parts
from meta import constants

from src.meta_utilities import capture_image,visualize_3d_critical_section
from src.general_utilities import add_row

class BOMF28DoorsSlide():

    def __init__(self,
                slide,
                windows,
                general_input,
                metadb_2d_input,
                metadb_3d_input,
                template_file,
                twod_images_report_folder,
                threed_images_report_folder,
                ppt_report_folder) -> None:
        self.shapes = slide.shapes
        self.windows = windows
        self.general_input = general_input
        self.metadb_2d_input = metadb_2d_input
        self.metadb_3d_input = metadb_3d_input
        self.template_file = template_file
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.ppt_report_folder = ppt_report_folder
        self.visible_parts = None

    def setup(self):
        """
        setup _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """
        utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
        utils.MetaCommand('0:options state original')
        utils.MetaCommand('options fringebar off')
        data = self.metadb_3d_input.critical_sections["f28_front_door"]
        visualize_3d_critical_section(data)
        m = models.Model(0)
        self.visible_parts = m.get_parts('visible')
        print("self.visible_parts",len(self.visible_parts))

        return 0

    def edit(self):
        """
        edit _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """
        from pptx.util import Pt

        self.setup()

        for shape in self.shapes:
            if shape.name == "Image 1":
                image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"f21_front_floor".lower()+".png")
                capture_image(self.general_input.threed_window_name,shape.width,shape.height,image_path)
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
            elif shape.name == "Table 1":
                table_obj = shape.table
                for index,prop in enumerate(self.visible_parts[:16]):
                    part = parts.Part(id=prop.id, type=constants.PSHELL, model_id=0)
                    materials = part.get_materials('all')

                    add_row(table_obj)
                    prop_row = table_obj.rows[index+1]

                    text_frame = prop_row.cells[0].text_frame
                    font = text_frame.paragraphs[0].font
                    font.size = Pt(8)
                    text_frame.paragraphs[0].text = str(prop.id)

                    text_frame_name = prop_row.cells[1].text_frame
                    font_name = text_frame_name.paragraphs[0].font
                    font_name.size = Pt(8)
                    text_frame_name.paragraphs[0].text = str(prop.name)

                    text_frame_material = prop_row.cells[2].text_frame
                    font_material = text_frame_material.paragraphs[0].font
                    font_material.size = Pt(8)
                    text_frame_material.paragraphs[0].text = str(materials[0].name)

                    text_frame_thickness = prop_row.cells[3].text_frame
                    font_thickness = text_frame_thickness.paragraphs[0].font
                    font_thickness.size = Pt(8)
                    thickness = round(prop.shell_thick,1)
                    text_frame_thickness.paragraphs[0].text = str(thickness)
            elif shape.name == "Table 2":
                table_obj = shape.table
                for index,prop in enumerate(self.visible_parts[16:]):
                    part = parts.Part(id=prop.id, type=constants.PSHELL, model_id=0)
                    materials = part.get_materials('all')

                    add_row(table_obj)
                    prop_row = table_obj.rows[index+1]

                    text_frame = prop_row.cells[0].text_frame
                    font = text_frame.paragraphs[0].font
                    font.size = Pt(8)
                    text_frame.paragraphs[0].text = str(prop.id)

                    text_frame_name = prop_row.cells[1].text_frame
                    font_name = text_frame_name.paragraphs[0].font
                    font_name.size = Pt(8)
                    text_frame_name.paragraphs[0].text = str(prop.name)

                    text_frame_material = prop_row.cells[2].text_frame
                    font_material = text_frame_material.paragraphs[0].font
                    font_material.size = Pt(8)
                    text_frame_material.paragraphs[0].text = str(materials[0].name)

                    text_frame_thickness = prop_row.cells[3].text_frame
                    font_thickness = text_frame_thickness.paragraphs[0].font
                    font_thickness.size = Pt(8)
                    thickness = round(prop.shell_thick,1)
                    text_frame_thickness.paragraphs[0].text = str(thickness)
        self.revert()

        return 0

    def revert(self):
        """
        revert _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """
        utils.MetaCommand('color pid transparency reset act')

        return 0
