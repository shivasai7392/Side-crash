# PYTHON script
"""
This script is used for all the automation process of Bill of materials f21 upb slide of thesis report.
"""

import os
import logging
from datetime import datetime

from meta import utils
from meta import parts
from meta import constants
from meta import models

from src.meta_utilities import capture_image
from src.meta_utilities import visualize_3d_critical_section
from src.general_utilities import add_row

class BOMF21UPBSlide():
    """
       This class is used to automate the BOM F21 upb slide of thesis report.

        Args:
            slide (object): bom f21 upb pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
            metadb_3d_input (Meta3DInfo): Meta3DInfo class object.
            threed_images_report_folder (str): folder path to save threed data images.
        """
    def __init__(self,
                slide,
                general_input,
                metadb_3d_input,
                threed_images_report_folder) -> None:
        self.shapes = slide.shapes
        self.general_input = general_input
        self.metadb_3d_input = metadb_3d_input
        self.threed_images_report_folder = threed_images_report_folder
        self.logger = logging.getLogger("side_crash_logger")

    def edit(self):
        """
        This method is used to iterate all the shapes of the bom f21 upb slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """
        from pptx.util import Pt
        try:
            self.logger.info("Started seeding data into bom f21 upb slide")
            self.logger.info("")
            starttime = datetime.now()
            m = models.Model(0)
            utils.MetaCommand('0:options state original')
            #iterating through the Image shapes of the bom f21 upb slide
            image_shapes = [shape for shape in self.shapes if "Image" in shape.name]
            for shape in image_shapes:
                #image insertion for the shape named "Image 2"
                if shape.name == "Image 2":
                    #visualizing "f21_upb_inner" critical part set to capture image at original state
                    data = self.metadb_3d_input.critical_sections["f21_upb_inner"]
                    visualize_3d_critical_section(data)
                    self.f21_upb_inner_visible_parts = m.get_parts('visible')
                    utils.MetaCommand('window maximize "{}"'.format(self.general_input.threed_window_name))
                    utils.MetaCommand('options fringebar off')
                    image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"F21_UPB_INNER"+".png").replace(" ","_")
                    capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height,view = "right",transparent=True)
                    self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                    self.logger.info("SOURCE MODEL : 0")
                    self.logger.info("STATE : ORIGINAL STATE")
                    self.logger.info("PID NAME SHOW FILTER : {} ".format(data["hes"] if "hes" in data.keys() else "null"))
                    self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                    self.logger.info("PID NAME ERASE FILTER : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                    self.logger.info("PID'S TO ERASE : {} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null"))
                    self.logger.info("ERASE BOX : {} ".format(data["erase_box"] if "erase_box" in data.keys() else "null"))
                    self.logger.info("IMAGE VIEW : {} ".format(data["view"] if "view" in data.keys() else "null"))
                    self.logger.info("TRANSPARENCY LEVEL : 50" )
                    self.logger.info("TRANSPARENT PID'S : {} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null"))
                    self.logger.info("COMP NAME : {} ".format(data["name"] if "name" in data.keys() else "null"))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT MODEL IMAGES :")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    transparent_image_path = image_path.replace(".png","")+"_transparent.png"
                    picture = self.shapes.add_picture(transparent_image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                    #removing transparent image
                    os.remove(transparent_image_path)
                    utils.MetaCommand('color pid transparency reset act')
                #image insertion for the shape named "Image 1"
                if shape.name == "Image 1":
                    #visualizing "f21_upb_outer" critical part set to capture image at original state
                    data = self.metadb_3d_input.critical_sections["f21_upb_outer"]
                    visualize_3d_critical_section(data)
                    self.f21_upb_outer_visible_parts = m.get_parts('visible')
                    utils.MetaCommand('window maximize "{}"'.format(self.general_input.threed_window_name))
                    utils.MetaCommand('options fringebar off')
                    image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"F21_UPB_OUTER"+".png").replace(" ","_")
                    capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height,view = "left",transparent=True)
                    self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                    self.logger.info("SOURCE MODEL : 0")
                    self.logger.info("STATE : ORIGINAL STATE")
                    self.logger.info("PID NAME SHOW FILTER : {} ".format(data["hes"] if "hes" in data.keys() else "null"))
                    self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                    self.logger.info("PID NAME ERASE FILTER : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                    self.logger.info("PID'S TO ERASE : {} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null"))
                    self.logger.info("ERASE BOX : {} ".format(data["erase_box"] if "erase_box" in data.keys() else "null"))
                    self.logger.info("IMAGE VIEW : {} ".format(data["view"] if "view" in data.keys() else "null"))
                    self.logger.info("TRANSPARENCY LEVEL : 50" )
                    self.logger.info("TRANSPARENT PID'S : {} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null"))
                    self.logger.info("COMP NAME : {} ".format(data["name"] if "name" in data.keys() else "null"))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT MODEL IMAGES :")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    transparent_image_path = image_path.replace(".png","")+"_transparent.png"
                    picture = self.shapes.add_picture(transparent_image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                    #removing transparent image
                    os.remove(transparent_image_path)
                    utils.MetaCommand('color pid transparency reset act')
            #iterating through the Table and Textbox shapes of the bom f21 upb slide
            remaining_shapes = [shape for shape in self.shapes if "Image" not in shape.name]
            for shape in remaining_shapes:
                #table population for the shape named "Table 1"
                if shape.name == "Table 1":
                    #getting the table object
                    table_obj = shape.table
                    #iterating through first 15 f21 upb outer visible parts
                    for index,prop in enumerate(self.f21_upb_outer_visible_parts[:15]):
                        #getting material object for the current part
                        part_type = parts.StringPartType(prop.type)
                        if part_type == "PSHELL":
                            part = parts.Part(id=prop.id,type = constants.PSHELL, model_id=0)
                            part_material = part.get_materials('all')[0]
                        elif part_type == "PSOLID":
                            part = parts.Part(id=prop.id,type = constants.PSOLID, model_id=0)
                            part_material = part.get_materials('all')[0]
                        #adding new row to the table
                        add_row(table_obj)
                        #getting the added row object
                        prop_row = table_obj.rows[index+1]
                        #getting cell 0 text frame object to insert property id
                        text_frame = prop_row.cells[0].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(8)
                        text_frame.paragraphs[0].text = str(prop.id)
                        #getting cell 1 text frame object to insert property name
                        text_frame_name = prop_row.cells[1].text_frame
                        font_name = text_frame_name.paragraphs[0].font
                        font_name.size = Pt(8)
                        text_frame_name.paragraphs[0].text = str(prop.name)
                        #getting cell 2 text frame object to insert property material name
                        text_frame_material = prop_row.cells[2].text_frame
                        font_material = text_frame_material.paragraphs[0].font
                        font_material.size = Pt(8)
                        text_frame_material.paragraphs[0].text = str(part_material.name)
                        #getting cell 3 text frame object to insert property shell thickness
                        text_frame_thickness = prop_row.cells[3].text_frame
                        font_thickness = text_frame_thickness.paragraphs[0].font
                        font_thickness.size = Pt(8)
                        thickness = round(prop.shell_thick,1)
                        text_frame_thickness.paragraphs[0].text = str(thickness)
                #table population for the shape named "Table 2"
                elif shape.name == "Table 2":
                    #getting the table object
                    table_obj = shape.table
                    #iterating through remaing f21 upb outer visible parts
                    for index,prop in enumerate(self.f21_upb_outer_visible_parts[15:]):
                        #getting material object for the current part
                        part_type = parts.StringPartType(prop.type)
                        if part_type == "PSHELL":
                            part = parts.Part(id=prop.id,type = constants.PSHELL, model_id=0)
                            part_material = part.get_materials('all')[0]
                        elif part_type == "PSOLID":
                            part = parts.Part(id=prop.id,type = constants.PSOLID, model_id=0)
                            part_material = part.get_materials('all')[0]
                        #adding new row to the table
                        add_row(table_obj)
                        #getting the added row object
                        prop_row = table_obj.rows[index+1]
                        #getting cell 0 text frame object to insert property id
                        text_frame = prop_row.cells[0].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(8)
                        text_frame.paragraphs[0].text = str(prop.id)
                        #getting cell 1 text frame object to insert property name
                        text_frame_name = prop_row.cells[1].text_frame
                        font_name = text_frame_name.paragraphs[0].font
                        font_name.size = Pt(8)
                        text_frame_name.paragraphs[0].text = str(prop.name)
                        #getting cell 2 text frame object to insert property material name
                        text_frame_material = prop_row.cells[2].text_frame
                        font_material = text_frame_material.paragraphs[0].font
                        font_material.size = Pt(8)
                        text_frame_material.paragraphs[0].text = str(part_material.name)
                        #getting cell 3 text frame object to insert property shell thickness
                        text_frame_thickness = prop_row.cells[3].text_frame
                        font_thickness = text_frame_thickness.paragraphs[0].font
                        font_thickness.size = Pt(8)
                        thickness = round(prop.shell_thick,1)
                        text_frame_thickness.paragraphs[0].text = str(thickness)
                #table population for the shape named "Table 3"
                elif shape.name == "Table 3":
                    #getting the table object
                    table_obj = shape.table
                    #iterating through remaing f21 upb inner visible parts
                    for index,prop in enumerate(self.f21_upb_inner_visible_parts):
                        #getting material object for the current part
                        part_type = parts.StringPartType(prop.type)
                        if part_type == "PSHELL":
                            part = parts.Part(id=prop.id,type = constants.PSHELL, model_id=0)
                            part_material = part.get_materials('all')[0]
                        elif part_type == "PSOLID":
                            part = parts.Part(id=prop.id,type = constants.PSOLID, model_id=0)
                            part_material = part.get_materials('all')[0]
                        #adding new row to the table
                        add_row(table_obj)
                        #getting the added row object
                        prop_row = table_obj.rows[index+1]
                        #getting cell 0 text frame object to insert property id
                        text_frame = prop_row.cells[0].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(8)
                        text_frame.paragraphs[0].text = str(prop.id)
                        #getting cell 1 text frame object to insert property name
                        text_frame_name = prop_row.cells[1].text_frame
                        font_name = text_frame_name.paragraphs[0].font
                        font_name.size = Pt(8)
                        text_frame_name.paragraphs[0].text = str(prop.name)
                        #getting cell 2 text frame object to insert property material name
                        text_frame_material = prop_row.cells[2].text_frame
                        font_material = text_frame_material.paragraphs[0].font
                        font_material.size = Pt(8)
                        text_frame_material.paragraphs[0].text = str(part_material.name)
                        #getting cell 3 text frame object to insert property shell thickness
                        text_frame_thickness = prop_row.cells[3].text_frame
                        font_thickness = text_frame_thickness.paragraphs[0].font
                        font_thickness.size = Pt(8)
                        thickness = round(prop.shell_thick,1)
                        text_frame_thickness.paragraphs[0].text = str(thickness)
                #text insertion for the shape named "TextBox 1"
                elif shape.name == "TextBox 1":
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "OUTER"
                #text insertion for the shape named "TextBox 2"
                elif shape.name == "TextBox 2":
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "INNER"
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into bom f21 upb slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into bom f21 upb slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")

        return 0
