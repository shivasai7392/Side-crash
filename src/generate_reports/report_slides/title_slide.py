# PYTHON script
"""
This script is used for all the automation process of Title slide of thesis report.
"""

import logging
from datetime import datetime

from src.metadb_info import GeneralVarInfo

class TitleSlide():
    """
        This class is used to automate the Title slide of thesis report.

        Args:
            slide (object): title pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
        """

    def __init__(self,
                slide,
                general_input) -> None:

        self.slide = slide
        self.shapes = slide.shapes
        self.general_input = general_input
        self.logger = logging.getLogger("side_crash_logger")

    def edit(self):
        """
        This method is used to iterate all the shapes of the title slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """
        from pptx.util import Pt

        try:
            self.logger.info("Started seeding data into title slide")
            self.logger.info("")
            self.logger.info(".")
            self.logger.info(".")
            self.logger.info(".")
            starttime = datetime.now()
            #iterating through the shapes of the title slide
            for shape in self.shapes:
                #table population for the shape named "Table 2"
                if shape.name == "Table 2":
                    #getting table object and rows
                    table_obj = shape.table
                    rows = table_obj.rows
                    #ordinal for the day number
                    ordinal = lambda n: "%d%s" % (n,"tsnrhtdd"[(n//10%10!=1)*(n%10<4)*n%10::4])
                    #getting current month,day and year
                    month = datetime.today().strftime('%B')
                    day = int(datetime.today().strftime('%d'))
                    year = datetime.today().strftime('%Y')
                    #inserting date in row 3 cell 0
                    text_frame_1 = rows[3].cells[0].text_frame
                    font = text_frame_1.paragraphs[0].font
                    font.size = Pt(16)
                    text_frame_1.paragraphs[0].text = " {} {}, {}".format(month,ordinal(day),year)
                    #inserting verification in row 2 cell 0
                    text_frame_2 = rows[2].cells[0].text_frame
                    font = text_frame_2.paragraphs[0].font
                    font.size = Pt(16)
                    if self.general_input.verification_mode not in ["null","none",""]:
                        text_frame_2.paragraphs[0].text = " " + self.general_input.verification_mode
                    else:
                        self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.verification_mode_key))
                        self.logger.info("")
                    #inserting run directory path in row 4 cell 2
                    text_frame_3 = rows[4].cells[2].text_frame
                    font = text_frame_3.paragraphs[0].font
                    font.size = Pt(16)
                    if self.general_input.run_directory.startswith("/"):
                        text_frame_3.paragraphs[0].text = " " + self.general_input.run_directory
                    else:
                        self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.run_directory_key))
                        self.logger.info("")
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into title slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into title slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")

        return 0
