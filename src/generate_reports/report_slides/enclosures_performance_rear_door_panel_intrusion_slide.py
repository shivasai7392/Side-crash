# PYTHON script
"""
This script is used for all the automation process of enclosures performance rear door panel intrusion slide of thesis report.
"""

import os
import logging
from datetime import datetime

from meta import utils
from meta import plot2d
from meta import windows

from src.meta_utilities import capture_image_and_resize
from src.general_utilities import add_row

class EnclosurePerformanceRearDoorPanelIntrusionSlide():
    """
       This class is used to automate the enclosure performance rear door panel intrusion slide of thesis report.

        Args:
            slide (object): enclosure performance rear door panel intrsusion pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
            twod_images_report_folder (str): folder path to save twod data images.
        """
    def __init__(self,
                slide,
                general_input,
                twod_images_report_folder) -> None:
        self.shapes = slide.shapes
        self.intrusion_areas = {"ROW 2":{}}
        self.general_input = general_input
        self.twod_images_report_folder = twod_images_report_folder
        self.logger = logging.getLogger("side_crash_logger")

    def intrusion_curve_format(self,source_window,curve,temporary_window_name,curve_name,target_curve = None):
        """
        This method is used for formatting axis,title options and attributes of curve of source window.

        Args:
            source_window (str): source plot window name.
            curve (object): source curve object.
            temporary_window_name (str): window name.
            curve_name (str): Curve name to set.
            target_curve(object):intrusion target Curve Object

        Returns:
            int: 0 Always for Sucess.1 for Failure.
        """
        #moving the curve to a Temporary window
        if target_curve:
            utils.MetaCommand('xyplot curve copy "{}" {},{}'.format(source_window,curve.id,target_curve.id))
        else:
            utils.MetaCommand('xyplot curve copy "{}" {}'.format(source_window,curve.id))
        utils.MetaCommand('xyplot create "{}"'.format(temporary_window_name))
        if target_curve:
            utils.MetaCommand('xyplot curve paste "{}" 0 {},{}'.format(temporary_window_name,curve.id,target_curve.id))
        else:
            utils.MetaCommand('xyplot curve paste "{}" 0 {}'.format(temporary_window_name,curve.id))
        try:
            utils.MetaCommand('xyplot curve copy "{}" {}'.format(source_window,curve.id))
            utils.MetaCommand('xyplot create "{}"'.format(temporary_window_name))
            utils.MetaCommand('xyplot curve paste "{}" 0 {}'.format(temporary_window_name,curve.id))
            win = windows.Window(temporary_window_name, page_id=0)
            curve = win.get_curves('all')[0]
            y_values = []
            for x in [0.03,0.04,0.05,0.06,0.07,0.08]:
                y_values.append(round(curve.get_y_values_from_x(specifier = 'first', xvalue =x)[0]))
            self.intrusion_areas[curve_name.rsplit(" ",1)[0]][curve_name.rsplit(" ",1)[1]] = y_values
            utils.MetaCommand('xyplot gridoptions line major style "{}" 0 2'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions yaxis active "{}" 0 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions labels yposition "{}" 0 left'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions labels yalign "{}" 0 left'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions axyrange "{}" 0 0 0 1200'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions axyrange "{}" 0 0 0 200'.format(temporary_window_name))
            utils.MetaCommand('xyplot gridoptions yspace "{}" 0 40'.format(temporary_window_name))
            utils.MetaCommand('xyplot plotoptions title set "{}" 0 "{}"'.format(temporary_window_name,curve_name))
            utils.MetaCommand('xyplot axisoptions ylabel set "{}" 0 "Intrusion [mm]"'.format(temporary_window_name))
            utils.MetaCommand('xyplot curve select "{}" all'.format(temporary_window_name))
            utils.MetaCommand('xyplot curve deselect "{}" {}'.format(temporary_window_name,target_curve.id))
            utils.MetaCommand('xyplot curve set style "{}" selected 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot curve set linewidth "{}" selected 9'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions ylabel font "{}" 0 "Arial,30,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions labels yfont "{}" 0 "Arial,30,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" 0 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions xaxis active "{}" 0 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions xlabel font "{}" 0 "Arial,30,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions labels xfont "{}" 0 "Arial,30,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot plotoptions title font "{}" 0 "Arial,30,-1,5,75,0,0,0,0,0"'.format(temporary_window_name))
            utils.MetaCommand('xyplot axisoptions xaxis deactive "{}" 0 0'.format(temporary_window_name))
            utils.MetaCommand('xyplot curve select "{}" all'.format(temporary_window_name))
        except:
            return 1

        return 0

    def edit(self):
        """
        This method is used to iterate all the shapes of the enclosure performance rear door panel intrusion slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """
        from pptx.util import Pt

        try:
            self.logger.info("Started seeding data into enclosure performance rear door panel intrusion slide")
            self.logger.info("")
            starttime = datetime.now()
            shapes = [shape for shape in self.shapes]
            shapes.sort(key = lambda x:x.name)
            #iterating through the Image shapes of the enclosure performance rear door panel intrusion slide
            survival_space_window_name = self.general_input.survival_space_window_name
            for shape in shapes:
                #image insertion for the shape named "Image 1"
                if shape.name == "Image 1":
                    #getting "Rear Door Accel" window and "Rear shoulder intrusion curve" curve objects to activate and format the curve axisoptions and title options
                    temporary_window_name = "Temporary"
                    rear_shoulder_intrusion_curve_name = self.general_input.rear_shoulder_intrusion_curve_name
                    rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                    utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                    curve = plot2d.CurvesByName(rear_door_accel_window_name, rear_shoulder_intrusion_curve_name, 1)[0]
                    target_curve = plot2d.CurvesByName(rear_door_accel_window_name, "*TARGET", 0)[0]
                    target_curve.show()
                    curve.show()
                    self.intrusion_curve_format(rear_door_accel_window_name,curve,temporary_window_name,"ROW 2 SHOULDER",target_curve=target_curve)
                    #capturing image of the formatted intrusion curve
                    image_path = os.path.join(self.twod_images_report_folder,survival_space_window_name+"_"+"ROW 2 SHOULDER"+".png").replace(" ","_")
                    capture_image_and_resize(image_path,shape.width,shape.height)
                    self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(self.general_input.rear_shoulder_intrusion_curve_name,self.general_input.rear_shoulder_intrusion_curve_key,rear_door_accel_window_name))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT CURVE IMAGES : ")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                    utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                #image insertion for the shape named "Image 2"
                elif shape.name == "Image 2":
                    #getting "Rear Door Accel" window and "Rear abdomen intrusion curve" curve objects to activate and format the curve axisoptions and title options
                    temporary_window_name = "Temporary"
                    rear_abdomen_intrusion_curve_name = self.general_input.rear_abdomen_intrusion_curve_name
                    rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                    utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                    curve = plot2d.CurvesByName(rear_door_accel_window_name, rear_abdomen_intrusion_curve_name, 1)[0]
                    target_curve = plot2d.CurvesByName(rear_door_accel_window_name, "*TARGET", 0)[0]
                    target_curve.show()
                    curve.show()
                    self.intrusion_curve_format(rear_door_accel_window_name,curve,temporary_window_name,"ROW 2 ABDOMEN",target_curve=target_curve)
                    #capturing image of the formatted intrusion curve
                    image_path = os.path.join(self.twod_images_report_folder,survival_space_window_name+"_"+"ROW 2 ABDOMEN"+".png").replace(" ","_")
                    capture_image_and_resize(image_path,shape.width,shape.height)
                    self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(self.general_input.rear_abdomen_intrusion_curve_name,self.general_input.rear_abdomen_intrusion_curve_key,rear_door_accel_window_name))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT CURVE IMAGES : ")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                    utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                #image insertion for the shape named "Image 2"
                elif shape.name == "Image 3":
                    #getting "Rear Door Accel" window and "Rear abdomen intrusion curve" curve objects to activate and format the curve axisoptions and title options
                    temporary_window_name = "Temporary"
                    rear_pelvis_intrusion_curve_name = self.general_input.rear_pelvis_intrusion_curve_name
                    rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                    utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                    curve = plot2d.CurvesByName(rear_door_accel_window_name, rear_pelvis_intrusion_curve_name, 1)[0]
                    target_curve = plot2d.CurvesByName(rear_door_accel_window_name, "*TARGET", 0)[0]
                    target_curve.show()
                    curve.show()
                    self.intrusion_curve_format(rear_door_accel_window_name,curve,temporary_window_name,"ROW 2 PELVIS",target_curve=target_curve)
                    #capturing image of the formatted intrusion curve
                    image_path = os.path.join(self.twod_images_report_folder,survival_space_window_name+"_"+"ROW 2 PELVIS"+".png").replace(" ","_")
                    capture_image_and_resize(image_path,shape.width,shape.height)
                    self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(self.general_input.rear_pelvis_intrusion_curve_name,self.general_input.rear_pelvis_intrusion_curve_key,rear_door_accel_window_name))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT CURVE IMAGES : ")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                    utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
                #image insertion for the shape named "Image 4"
                elif shape.name == "Image 4":
                    #getting "Rear Door Accel" window and "Rear femur intrusion curve" curve objects to activate and format the curve axisoptions and title options
                    temporary_window_name = "Temporary"
                    rear_femur_intrusion_curve_name = self.general_input.rear_femur_intrusion_curve_name
                    rear_door_accel_window_name = self.general_input.rear_door_accel_window_name
                    utils.MetaCommand('window maximize {}'.format(rear_door_accel_window_name))
                    curve = plot2d.CurvesByName(rear_door_accel_window_name, rear_femur_intrusion_curve_name, 1)[0]
                    target_curve = plot2d.CurvesByName(rear_door_accel_window_name, "*TARGET", 0)[0]
                    target_curve.show()
                    curve.show()
                    self.intrusion_curve_format(rear_door_accel_window_name,curve,temporary_window_name,"ROW 2 FEMUR",target_curve=target_curve)
                    #capturing image of the formatted intrusion curve
                    image_path = os.path.join(self.twod_images_report_folder,survival_space_window_name+"_"+"ROW 2 FEMUR"+".png").replace(" ","_")
                    capture_image_and_resize(image_path,shape.width,shape.height)
                    self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("CURVES : {} | FROM VARIABLE : {} | SOURCE WINDOW : {}".format(self.general_input.rear_femur_intrusion_curve_name,self.general_input.rear_femur_intrusion_curve_key,rear_door_accel_window_name))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT CURVE IMAGES : ")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                    utils.MetaCommand('window delete "{}"'.format(temporary_window_name))
            #iterating through the Table shapes of the enclosure performance rear door panel intrusion slide
            for shape in self.shapes:
                #table population for the table named as "Table 1"
                if shape.name == "Table 1":
                    #getting table object
                    table = shape.table
                    row_index = 2
                    node_index = 0
                    #adding new row to the table
                    add_row(table)
                    #iterating through intrusion curve data of "ROW 2"
                    for iindex,(key,values) in enumerate(self.intrusion_areas["ROW 2"].items()):
                        if iindex == 2:
                            #adding a new row
                            add_row(table)
                            row_index = row_index+1
                        if iindex == 1 or iindex == 3:
                            node_index = 7
                        else:
                            node_index = 0
                        #getting table row objects
                        rows = table.rows
                        #inserting heading into the cell
                        text_frame_1 = rows[row_index].cells[node_index].text_frame
                        font = text_frame_1.paragraphs[0].font
                        font.bold = True
                        font.size = Pt(11)
                        text_frame_1.paragraphs[0].text = key.capitalize()
                        #inserting y values in to the cells
                        for index,value in enumerate(values):
                            text_frame = rows[row_index].cells[node_index+index+1].text_frame
                            font = text_frame.paragraphs[0].font
                            font.size = Pt(11)
                            text_frame.paragraphs[0].text = str(value)
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into enclosure performance rear door panel intrusion slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into enclosure performance rear door panel intrusion slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")

        return 0
