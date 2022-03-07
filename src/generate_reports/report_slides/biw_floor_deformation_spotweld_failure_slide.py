# PYTHON script
"""
    _summary_

_extended_summary_

Returns:
    _type_: _description_
"""

import os

from meta import utils
from meta import windows
from meta import plot2d

from src.meta_utilities import capture_image,visualize_3d_critical_section,capture_resized_image
from src.general_utilities import closest

class BIWFloorDeformationAndSpotWeldFailureSlide():

    def __init__(self,
                slide,
                meta_windows,
                general_input,
                metadb_2d_input,
                metadb_3d_input,
                template_file,
                twod_images_report_folder,
                threed_images_report_folder,
                ppt_report_folder) -> None:
        self.shapes = slide.shapes
        self.windows = meta_windows
        self.general_input = general_input
        self.metadb_2d_input = metadb_2d_input
        self.metadb_3d_input = metadb_3d_input
        self.template_file = template_file
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.ppt_report_folder = ppt_report_folder
        self.visible_parts = None

    def setup(self):
        """
        setup _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """

        return 0

    def edit(self):
        """
        edit _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """
        from pptx.util import Pt

        self.setup()

        for shape in self.shapes:
            if shape.name == "Image 3":
                utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('grstyle scalarfringe enable')
                utils.MetaCommand('options fringebar off')
                data = self.metadb_3d_input.critical_sections["f21_front_floor"]
                visualize_3d_critical_section(data)
                utils.MetaCommand('color pid transparency reset act')
                image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"f21_front_floor".lower()+".png")
                capture_image(self.general_input.threed_window_name,shape.width,shape.height,image_path,view = "btm")
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('options fringebar on')
                utils.MetaCommand('0:options state original"')
            elif shape.name == "Image 4":
                utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('grstyle scalarfringe enable')
                utils.MetaCommand('options fringebar off')
                utils.MetaCommand('grstyle deform off')
                data = self.metadb_3d_input.critical_sections["f21_front_floor"]
                visualize_3d_critical_section(data)
                utils.MetaCommand('color pid transparency reset act')
                image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"f21_front_floor".lower()+".png")
                capture_image(self.general_input.threed_window_name,shape.width,shape.height,image_path,view = "btm")
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('options fringebar on')
                utils.MetaCommand('grstyle deform on')
                utils.MetaCommand('0:options state original"')
            elif shape.name == "Image 2":
                utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('options fringebar on')
                image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"fringe_bar".lower()+".png")
                utils.MetaCommand('write scalarfringebar png {} '.format(image_path))
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('options fringebar off')
            elif shape.name == "Image 1":
                window_name = self.general_input.biw_stiff_ring_deformation_name
                win = windows.Window(str(window_name), page_id=0)
                layout = win.get_plot_layout()
                utils.MetaCommand('window maximize "{}"'.format(window_name))
                plot_id = 1
                page_id=0
                final_time_variable =  dict(utils.MetaGetVariablesByName("survival-space_final_time"))
                final_time_roof = final_time_variable["survival-space_final_time"]
                final_time_roof_splitting = final_time_roof.split(".")[0]
                plot = plot2d.Plot(plot_id, window_name, page_id)
                curvelist_final_time = plot.get_curves('byname', name ="SIDE_SILL_"+str(final_time_roof_splitting)+"MS")
                for each_curvelist_final_time in curvelist_final_time:
                    final_time_id = each_curvelist_final_time.id
                curvelist_initial_time = plot.get_curves('byname', name ="SIDE_SILL_"+str(0)+"MS")

                peak_time_variable = dict(utils.MetaGetVariablesByName("peak_time_display"))
                peak_time_value = peak_time_variable["peak_time_display"]
                peak_time = peak_time_value.split(".")[0]
                plot = plot2d.Plot(plot_id, window_name, page_id)
                curves = plot.get_curves('all')
                final_roof_line_list = list()
                for each_curve in curves:
                    ms = each_curve.name.split("_")[2]
                    if 'MS' in ms:
                        ms_replacing = ms.replace('MS',"")
                        final_roof_line_list.append(int(ms_replacing))
                peak_time_value = closest(final_roof_line_list, int(peak_time))
                peak_time_curve = plot.get_curves('byname', name ="SIDE_SILL_"+str(peak_time_value)+"MS")
                for each_peak_time_curve in peak_time_curve:
                    peak_time_id = each_peak_time_curve.id

                for each_curvelist_initial_time in curvelist_initial_time:
                    initial_time_id = each_curvelist_initial_time.id

                title = plot2d.Title(plot_id, window_name, page_id)
                plot = title.get_plot()
                plot.activate()
                utils.MetaCommand('xyplot plotactive "{}" 1'.format(window_name))
                utils.MetaCommand('xyplot rlayout "{}" 1'.format(window_name))
                utils.MetaCommand('xyplot curve visible and "{}" {} {},{}'.format(window_name,initial_time_id,peak_time_id, final_time_id))
                utils.MetaCommand('xyplot curve set style "{}" {} 9'.format(window_name, initial_time_id))
                utils.MetaCommand('xyplot curve set style "{}" {} 5'.format(window_name,peak_time_id))
                utils.MetaCommand('xyplot curve select "{}" all'.format(window_name))
                utils.MetaCommand('xyplot axisoptions yaxis active "{}" 1 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions axyrange "{}" 1 0 -805 -770'.format(window_name))
                utils.MetaCommand('xyplot axisoptions ylabel font "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions labels yfont "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions yaxis deactive "{}" 1 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions xaxis active "{}" 1 0'.format(window_name))
                utils.MetaCommand('xyplot axisoptions xlabel font "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot axisoptions labels xfont "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))
                utils.MetaCommand('xyplot plotoptions title font "{}" 1 "Arial,32,-1,5,75,0,0,0,0,0"'.format(window_name))

                image_path = os.path.join(self.twod_images_report_folder,window_name+"_"+title.get_text().lower()+".png").replace("&","and")
                capture_resized_image(window_name,shape.width,shape.height,image_path,plot_id=plot.id)
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                plot.deactivate()
                utils.MetaCommand('xyplot rlayout "{}" {}'.format(window_name, layout))

        self.revert()

        return 0

    def revert(self):
        """
        revert _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """

        return 0