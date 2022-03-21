# PYTHON script
"""
This script is used for all the automation process of Enclosure performance skin defirmation slide of thesis report.
"""

import os
import logging
from datetime import datetime

from meta import utils
from meta import windows

from src.meta_utilities import capture_image
from src.meta_utilities import visualize_3d_critical_section
from src.meta_utilities import capture_image_and_resize
from src.general_utilities import closest

class EnclosurePerformanceSkinDeformationSlide():
    """
       This class is used to automate the enclosure performance skin deformation slide of thesis report.

        Args:
            slide (object): enclosure performance skin deformation pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
            metadb_3d_input (Meta3DInfo): Meta3DInfo class object.
            twod_images_report_folder (str): folder path to save twod data images.
            threed_images_report_folder (str): folder path to save threed data images.
        """
    def __init__(self,
                slide,
                general_input,
                metadb_3d_input,
                twod_images_report_folder,
                threed_images_report_folder) -> None:
        self.shapes = slide.shapes
        self.general_input = general_input
        self.metadb_3d_input = metadb_3d_input
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.visible_parts = None
        self.logger = logging.getLogger("side_crash_logger")

    def edit(self):
        """
        This method is used to iterate all the shapes of the enclosure performance skin deformation slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """
        from pptx.util import Pt

        try:
            self.logger.info("Started seeding data into enclosure performance skin deformation slide")
            self.logger.info("")
            starttime = datetime.now()
            #iterating through the Image shapes of the enclosure performance skin deformation slide
            for shape in self.shapes:
                #image insertion for the shape named "Image 1"
                if shape.name == "Image 1":
                    #capturing fringe bar of metapost window
                    utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                    utils.MetaCommand('0:options state variable "serial=1"')
                    utils.MetaCommand('options fringebar on')
                    image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"fringe_bar".lower()+".png")
                    utils.MetaCommand('write scalarfringebar png {} '.format(image_path))
                    utils.MetaCommand('write scalarfringebar png {} '.format(image_path))
                    self.logger.info("--- 3D FRINGE BAR IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT MODEL IMAGES :")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                    utils.MetaCommand('options fringebar off')
                #image insertion for the shape named "Image 2"
                elif shape.name == "Image 2":
                    #visualizing and capturing image of "f28_front_door and f28_rear_door" critical part sets at peak state with deformation
                    utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                    data = self.metadb_3d_input.critical_sections["f28_front_door"]
                    visualize_3d_critical_section(data)
                    data2 = self.metadb_3d_input.critical_sections["f28_rear_door"]
                    visualize_3d_critical_section(data2,and_filter = True)
                    utils.MetaCommand('0:options state variable "serial=1"')
                    utils.MetaCommand('color pid transparency reset act')
                    utils.MetaCommand('grstyle scalarfringe enable')
                    utils.MetaCommand('options fringebar off')
                    utils.MetaCommand('grstyle deform on')
                    image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"F21_FRONT_DOOR_AND_REAR_DOOR_AT_PEAK_STATE_WITH_DEFORMTION"+".jpeg")
                    capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height)
                    self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                    self.logger.info("SOURCE MODEL : 0")
                    self.logger.info("STATE : ORIGINAL STATE")
                    self.logger.info("PID NAME SHOW FILTER : {},{} ".format(data["hes"] if "hes" in data.keys() else "null",data2["hes"] if "hes" in data2.keys() else "null"))
                    self.logger.info("ADDITIONAL PID'S SHOWN : {},{} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null",data2["hes_exceptions"] if "hes_exceptions" in data2.keys() else "null"))
                    self.logger.info("PID NAME ERASE FILTER : {},{} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null",data2["hes_exceptions"] if "hes_exceptions" in data2.keys() else "null"))
                    self.logger.info("PID'S TO ERASE : {},{} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null",data2["erase_pids"] if "erase_pids" in data2.keys() else "null"))
                    self.logger.info("ERASE BOX : {},{} ".format(data["erase_box"] if "erase_box" in data.keys() else "null",data2["erase_box"] if "erase_box" in data2.keys() else "null"))
                    self.logger.info("IMAGE VIEW : {},{} ".format(data["view"] if "view" in data.keys() else "null",data2["view"] if "view" in data2.keys() else "null"))
                    self.logger.info("TRANSPARENCY LEVEL : 50" )
                    self.logger.info("TRANSPARENT PID'S : {},{} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null",data2["transparent_pids"] if "transparent_pids" in data2.keys() else "null"))
                    self.logger.info("COMP NAME : {},{} ".format(data["name"] if "name" in data.keys() else "null",data2["name"] if "name" in data2.keys() else "null"))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT MODEL IMAGES :")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                #image insertion for the shape named "Image 3"
                elif shape.name == "Image 3":
                    #visualizing and capturing image of "f28_front_door and f28_rear_door" critical part sets at peak state without deformation
                    utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                    data = self.metadb_3d_input.critical_sections["f28_front_door"]
                    visualize_3d_critical_section(data)
                    data2 = self.metadb_3d_input.critical_sections["f28_rear_door"]
                    visualize_3d_critical_section(data2,and_filter = True)
                    utils.MetaCommand('0:options state variable "serial=1"')
                    utils.MetaCommand('color pid transparency reset act')
                    utils.MetaCommand('grstyle scalarfringe enable')
                    utils.MetaCommand('options fringebar off')
                    utils.MetaCommand('grstyle deform off')
                    image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"F21_FRONT_DOOR_AND_REAR_DOOR_AT_PEAK_STATE_WITHOUT_DEFORMTION"+".jpeg")
                    capture_image(image_path,self.general_input.threed_window_name,shape.width,shape.height)
                    self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                    self.logger.info("SOURCE MODEL : 0")
                    self.logger.info("STATE : ORIGINAL STATE")
                    self.logger.info("PID NAME SHOW FILTER : {},{} ".format(data["hes"] if "hes" in data.keys() else "null",data2["hes"] if "hes" in data2.keys() else "null"))
                    self.logger.info("ADDITIONAL PID'S SHOWN : {},{} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null",data2["hes_exceptions"] if "hes_exceptions" in data2.keys() else "null"))
                    self.logger.info("PID NAME ERASE FILTER : {},{} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null",data2["hes_exceptions"] if "hes_exceptions" in data2.keys() else "null"))
                    self.logger.info("PID'S TO ERASE : {},{} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null",data2["erase_pids"] if "erase_pids" in data2.keys() else "null"))
                    self.logger.info("ERASE BOX : {},{} ".format(data["erase_box"] if "erase_box" in data.keys() else "null",data2["erase_box"] if "erase_box" in data2.keys() else "null"))
                    self.logger.info("IMAGE VIEW : {},{} ".format(data["view"] if "view" in data.keys() else "null",data2["view"] if "view" in data2.keys() else "null"))
                    self.logger.info("TRANSPARENCY LEVEL : 50" )
                    self.logger.info("TRANSPARENT PID'S : {},{} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null",data2["transparent_pids"] if "transparent_pids" in data2.keys() else "null"))
                    self.logger.info("COMP NAME : {},{} ".format(data["name"] if "name" in data.keys() else "null",data2["name"] if "name" in data2.keys() else "null"))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT MODEL IMAGES :")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                    utils.MetaCommand('grstyle deform on')
                #image insertion for the shape named "Image 4"
                elif shape.name == "Image 4":
                    #getting "Door Skin intrusion" visible plot objects to activate and showing initial,finaland peak time curves
                    door_skin_intrusion_window_name = self.general_input.door_skin_intrusion_window_name
                    utils.MetaCommand('window maximize {}'.format(door_skin_intrusion_window_name))
                    door_skin_intrusion_window = windows.Window(door_skin_intrusion_window_name, page_id=0)
                    door_skin_intrusion_window_plots = door_skin_intrusion_window.get_plots('visible')
                    for plot in door_skin_intrusion_window_plots:
                        plot.activate()
                        initial_curve_re  ="*0MS"
                        final_curve_re = "*150MS"
                        curves = plot.get_curves('all')
                        skin_intrusion_line_list = []
                        for each_curve in curves:
                            ms = each_curve.name.rsplit("_",1)[1]
                            if 'MS' in ms:
                                ms_replacing = ms.replace('MS',"")
                                skin_intrusion_line_list.append(int(ms_replacing))
                        peak_time = str(closest(skin_intrusion_line_list, round(float(self.general_input.peak_time_display_value))))
                        peak_curve_re = "*{}MS".format(peak_time)
                        for name in [initial_curve_re,final_curve_re,peak_curve_re]:
                            curve = plot.get_curves('byname', name = name)[0]
                            curve.show()
                        curve.set_line_style(line_style = 5)
                        plot.deactivate()
                    #capturing the door skin intrusion window
                    image_path = os.path.join(self.twod_images_report_folder,door_skin_intrusion_window_name.lower()+".png")
                    capture_image_and_resize(image_path,shape.width,shape.height)
                    self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("CURVES : {} | SOURCE PLOT : VISIBLE PLOTS | SOURCE WINDOW : {}".format("*0MS,*150MS,*{}MS".format(peak_time),door_skin_intrusion_window_name))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT CURVE IMAGES : ")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into enclosure performance skin deformation slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into enclosure performance skin deformation slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")

        return 0
