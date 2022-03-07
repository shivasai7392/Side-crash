# PYTHON script
"""
    _summary_

_extended_summary_

Returns:
    _type_: _description_
"""

import os

from meta import utils
from meta import plot2d
from meta import windows

from src.meta_utilities import capture_image,visualize_3d_critical_section,capture_resized_image
from src.general_utilities import closest

class EnclosurePerformanceRearDoorPanelDeformationSlide():

    def __init__(self,
                slide,
                meta_windows,
                general_input,
                metadb_2d_input,
                metadb_3d_input,
                template_file,
                twod_images_report_folder,
                threed_images_report_folder,
                ppt_report_folder) -> None:
        self.shapes = slide.shapes
        self.windows = meta_windows
        self.general_input = general_input
        self.metadb_2d_input = metadb_2d_input
        self.metadb_3d_input = metadb_3d_input
        self.template_file = template_file
        self.twod_images_report_folder = twod_images_report_folder
        self.threed_images_report_folder = threed_images_report_folder
        self.ppt_report_folder = ppt_report_folder
        self.visible_parts = None

    def setup(self):
        """
        setup _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """

        return 0

    def edit(self):
        """
        edit _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """
        from pptx.util import Pt

        self.setup()
        for shape in self.shapes:
            if shape.name == "Image 1":
                utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('options fringebar on')
                image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"fringe_bar".lower()+".png")
                utils.MetaCommand('write scalarfringebar png {} '.format(image_path))
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('options fringebar off')
            elif shape.name == "Image 2":
                utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                data = self.metadb_3d_input.critical_sections["f28_rear_door"]
                visualize_3d_critical_section(data)
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('grstyle scalarfringe enable')
                utils.MetaCommand('options fringebar off')
                utils.MetaCommand('grstyle deform on')
                image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"f28_rear_door".lower()+".png")
                capture_image(self.general_input.threed_window_name,shape.width,shape.height,image_path,view = "left")
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('color pid transparency reset act')
            elif shape.name == "Image 3":
                utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
                data = self.metadb_3d_input.critical_sections["f28_rear_door"]
                visualize_3d_critical_section(data,and_filter = True)
                utils.MetaCommand('0:options state variable "serial=1"')
                utils.MetaCommand('grstyle scalarfringe enable')
                utils.MetaCommand('options fringebar off')
                utils.MetaCommand('grstyle deform off')
                image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"f28_rear_door".lower()+".png")
                capture_image(self.general_input.threed_window_name,shape.width,shape.height,image_path,view = "left")
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('grstyle deform on')
                utils.MetaCommand('color pid transparency reset act')
            if shape.name == "Image 4":
                door_panel_intrusion_window_name = self.general_input.door_panel_intrusion_window_name
                utils.MetaCommand('window maximize {}'.format(door_panel_intrusion_window_name))
                door_panel_intrusion_window = windows.Window(door_panel_intrusion_window_name, page_id=0)
                door_panel_intrusion_window_layout = door_panel_intrusion_window.get_plot_layout()
                door_skin_intrusion_window_plots = [plot2d.Plot(plot_id, door_panel_intrusion_window_name, 0) for plot_id in [1,3,5]]
                for plot in door_skin_intrusion_window_plots:
                    plot.activate()
                    initial_curve_re  ="*0MS"
                    final_curve_re = "*150MS"
                    curves = plot.get_curves('all')
                    deformation_line_list = []
                    for each_curve in curves:
                        ms = each_curve.name.rsplit("_",1)[1]
                        if 'MS' in ms:
                            ms_replacing = ms.replace('MS',"")
                            deformation_line_list.append(int(ms_replacing))
                    peak_curve_re = "*{}MS".format(str(closest(deformation_line_list, round(float(self.general_input.peak_time_display_value)))))
                    for name in [initial_curve_re,final_curve_re,peak_curve_re]:
                        curve = plot.get_curves('byname', name = name)[0]
                        curve.show()
                    curve.set_line_style(line_style = 5)
                utils.MetaCommand('xyplot rlayout "Door panel intrusion" 6')
                image_path = os.path.join(self.twod_images_report_folder,door_panel_intrusion_window_name.lower()+"deformation"+".png")
                capture_resized_image(door_panel_intrusion_window_name,shape.width,shape.height,image_path)
                picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                picture.crop_left = 0
                picture.crop_right = 0
                utils.MetaCommand('xyplot rlayout "Door panel intrusion" {}'.format(door_panel_intrusion_window_layout))
                utils.MetaCommand('xyplot plotdeactive "{}" all'.format(door_panel_intrusion_window_name))

        self.revert()

        return 0

    def revert(self):
        """
        revert _summary_

        _extended_summary_

        Returns:
            _type_: _description_
        """

        return 0
