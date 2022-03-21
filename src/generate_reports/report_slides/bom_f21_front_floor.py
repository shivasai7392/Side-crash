# PYTHON script
"""
This script is used for all the automation process of Bill of materials f21 front floor slide of thesis report.
"""

import os
import logging
from datetime import datetime

from meta import utils
from meta import models
from meta import parts
from meta import constants

from src.meta_utilities import capture_image
from src.meta_utilities import visualize_3d_critical_section
from src.general_utilities import add_row

class BOMF21FrontFloorSlide():
    """
       This class is used to automate the  Bill of materials f21 front floor slide of thesis report.

        Args:
            slide (object): bom f21 front floor pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
            metadb_3d_input (Meta3DInfo): Meta3DInfo class object.
            threed_images_report_folder (str): folder path to save threed data images.
        """
    def __init__(self,
                slide,
                general_input,
                metadb_3d_input,
                threed_images_report_folder) -> None:
        self.shapes = slide.shapes
        self.general_input = general_input
        self.metadb_3d_input = metadb_3d_input
        self.threed_images_report_folder = threed_images_report_folder
        self.visible_parts = None
        self.logger = logging.getLogger("side_crash_logger")

    def edit(self):
        """
        This method is used to iterate all the shapes of the bom f21 front floor slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """
        from pptx.util import Pt

        try:
            self.logger.info("Started seeding data into bom f21 front floor slide")
            self.logger.info("")
            starttime = datetime.now()
            #visualizing "f21_front_floor" critical part set to capture image at original state
            utils.MetaCommand('window maximize {}'.format(self.general_input.threed_window_name))
            utils.MetaCommand('0:options state original')
            utils.MetaCommand('options fringebar off')
            data = self.metadb_3d_input.critical_sections["f21_front_floor"]
            visualize_3d_critical_section(data)
            m = models.Model(0)
            self.visible_parts = m.get_parts('visible')
            #iterating through the Image shapes of the bom f21 front floor
            for shape in self.shapes:
                #image insertion for the shape named "Image 1"
                if shape.name == "Image 1":
                    image_path = os.path.join(self.threed_images_report_folder,self.general_input.threed_window_name+"_"+"F21_FRONT_FLOOR"+".jpeg")
                    #capturing the "f21_front_floor" critical part set image at original state
                    capture_image(image_path,shape.width,shape.height)
                    self.logger.info("--- 3D MODEL IMAGE GENERATOR")
                    self.logger.info("")
                    self.logger.info("SOURCE WINDOW : {} ".format(self.general_input.threed_window_name))
                    self.logger.info("SOURCE MODEL : 0")
                    self.logger.info("STATE : ORIGINAL STATE")
                    self.logger.info("PID NAME SHOW FILTER : {} ".format(data["hes"] if "hes" in data.keys() else "null"))
                    self.logger.info("ADDITIONAL PID'S SHOWN : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                    self.logger.info("PID NAME ERASE FILTER : {} ".format(data["hes_exceptions"] if "hes_exceptions" in data.keys() else "null"))
                    self.logger.info("PID'S TO ERASE : {} ".format(data["erase_pids"] if "erase_pids" in data.keys() else "null"))
                    self.logger.info("ERASE BOX : {} ".format(data["erase_box"] if "erase_box" in data.keys() else "null"))
                    self.logger.info("IMAGE VIEW : {} ".format(data["view"] if "view" in data.keys() else "null"))
                    self.logger.info("TRANSPARENCY LEVEL : 50" )
                    self.logger.info("TRANSPARENT PID'S : {} ".format(data["transparent_pids"] if "transparent_pids" in data.keys() else "null"))
                    self.logger.info("COMP NAME : {} ".format(data["name"] if "name" in data.keys() else "null"))
                    self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                    self.logger.info("OUTPUT MODEL IMAGES :")
                    self.logger.info(image_path)
                    self.logger.info("")
                    #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                    picture = self.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                    picture.crop_left = 0
                    picture.crop_right = 0
                #table population for the shape named "Table 1"
                elif shape.name == "Table 1":
                    #getting the table object
                    table_obj = shape.table
                    #iterating through f21 front floor first 13 visible parts
                    for index,prop in enumerate(self.visible_parts[:13]):
                        part_type = parts.StringPartType(prop.type)
                        if part_type == "PSHELL":
                            part = parts.Part(id=prop.id,type = constants.PSHELL, model_id=0)
                            part_material = part.get_materials('all')[0]
                        elif part_type == "PSOLID":
                            part = parts.Part(id=prop.id,type = constants.PSOLID, model_id=0)
                            part_material = part.get_materials('all')[0]
                        #adding new row to the table
                        add_row(table_obj)
                        #getting the added row object
                        prop_row = table_obj.rows[index+1]
                        #getting cell 0 text frame object to insert property id
                        text_frame = prop_row.cells[0].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(8)
                        text_frame.paragraphs[0].text = str(prop.id)
                        #getting cell 1 text frame object to insert property name
                        text_frame_name = prop_row.cells[1].text_frame
                        font_name = text_frame_name.paragraphs[0].font
                        font_name.size = Pt(8)
                        text_frame_name.paragraphs[0].text = str(prop.name)
                        #getting cell 2 text frame object to insert property material name
                        text_frame_material = prop_row.cells[2].text_frame
                        font_material = text_frame_material.paragraphs[0].font
                        font_material.size = Pt(8)
                        text_frame_material.paragraphs[0].text = str(part_material.name)
                        #getting cell 3 text frame object to insert property shell thickness
                        text_frame_thickness = prop_row.cells[3].text_frame
                        font_thickness = text_frame_thickness.paragraphs[0].font
                        font_thickness.size = Pt(8)
                        thickness = round(prop.shell_thick,1)
                        text_frame_thickness.paragraphs[0].text = str(thickness)
                #table population for the shape named "Table 2"
                elif shape.name == "Table 2":
                    #getting the table object
                    table_obj = shape.table
                    #iterating through f21 front floor for next 12 visible parts
                    for index,prop in enumerate(self.visible_parts[13:27]):
                        part_type = parts.StringPartType(prop.type)
                        if part_type == "PSHELL":
                            part = parts.Part(id=prop.id,type = constants.PSHELL, model_id=0)
                            part_material = part.get_materials('all')[0]
                        elif part_type == "PSOLID":
                            part = parts.Part(id=prop.id,type = constants.PSOLID, model_id=0)
                            part_material = part.get_materials('all')[0]
                        #adding new row to the table
                        add_row(table_obj)
                        #getting the added row object
                        prop_row = table_obj.rows[index+1]
                        #getting cell 0 text frame object to insert property id
                        text_frame = prop_row.cells[0].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(8)
                        text_frame.paragraphs[0].text = str(prop.id)
                        #getting cell 1 text frame object to insert property id
                        text_frame_name = prop_row.cells[1].text_frame
                        font_name = text_frame_name.paragraphs[0].font
                        font_name.size = Pt(8)
                        text_frame_name.paragraphs[0].text = str(prop.name)
                        #getting cell 2 text frame object to insert property name
                        text_frame_material = prop_row.cells[2].text_frame
                        font_material = text_frame_material.paragraphs[0].font
                        font_material.size = Pt(8)
                        text_frame_material.paragraphs[0].text = str(part_material.name)
                        #getting cell 3 text frame object to insert property material name
                        text_frame_thickness = prop_row.cells[3].text_frame
                        font_thickness = text_frame_thickness.paragraphs[0].font
                        font_thickness.size = Pt(8)
                        thickness = round(prop.shell_thick,1)
                        text_frame_thickness.paragraphs[0].text = str(thickness)
                #table population for the shape named "Table 3"
                elif shape.name == "Table 3":
                    #getting the table object
                    table_obj = shape.table
                    #iterating through f21 front floor for remaining visible parts
                    for index,prop in enumerate(self.visible_parts[27:]):
                        part_type = parts.StringPartType(prop.type)
                        if part_type == "PSHELL":
                            part = parts.Part(id=prop.id,type = constants.PSHELL, model_id=0)
                            part_material = part.get_materials('all')[0]
                        elif part_type == "PSOLID":
                            part = parts.Part(id=prop.id,type = constants.PSOLID, model_id=0)
                            part_material = part.get_materials('all')[0]
                        #adding new row to the table
                        add_row(table_obj)
                        #getting the added row object
                        prop_row = table_obj.rows[index+1]
                        #getting cell 0 text frame object to insert property id
                        text_frame = prop_row.cells[0].text_frame
                        font = text_frame.paragraphs[0].font
                        font.size = Pt(8)
                        text_frame.paragraphs[0].text = str(prop.id)
                        #getting cell 1 text frame object to insert property name
                        text_frame_name = prop_row.cells[1].text_frame
                        font_name = text_frame_name.paragraphs[0].font
                        font_name.size = Pt(8)
                        text_frame_name.paragraphs[0].text = str(prop.name)
                        #getting cell 2 text frame object to insert property material name
                        text_frame_material = prop_row.cells[2].text_frame
                        font_material = text_frame_material.paragraphs[0].font
                        font_material.size = Pt(8)
                        text_frame_material.paragraphs[0].text = str(part_material.name)
                        #getting cell 3 text frame object to insert property material name
                        text_frame_thickness = prop_row.cells[3].text_frame
                        font_thickness = text_frame_thickness.paragraphs[0].font
                        font_thickness.size = Pt(8)
                        thickness = round(prop.shell_thick,1)
                        text_frame_thickness.paragraphs[0].text = str(thickness)
            #revertting visual settings
            utils.MetaCommand('color pid transparency reset act')
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into bom f21 front floor slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into bom f21 front floor slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")

        return 0
