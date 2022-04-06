# PYTHON script
"""
This script is used for all the automation process of CAE Quality slide of thesis report.
"""

import os
import logging
from datetime import datetime

from meta import utils
from meta import plot2d
from meta import windows

from src.general_utilities import add_row
from src.meta_utilities import capture_image_and_resize
from src.metadb_info import GeneralVarInfo

class CAEQualitySlide():
    """
       This class is used to automate the CAE Quality slide of thesis report.

        Args:
            slide (object): cae quality pptx slide object.
            general_input (GeneralInfo): GeneralInfo class object.
            twod_images_report_folder (str): folder path to save twod data images.
        """

    def __init__(self,
                slide,
                general_input,
                twod_images_report_folder) -> None:
        self.slide = slide
        self.shapes = slide.shapes
        self.general_input = general_input
        self.twod_images_report_folder = twod_images_report_folder
        self.logger = logging.getLogger("side_crash_logger")

    def edit(self):
        """
        This method is used to iterate all the shapes of the cae quality slide and insert respective data.

        Returns:
            int: 0 Always for Sucess,1 for Failure.
        """

        from PIL import Image
        from PIL import ImageFile
        from pptx.util import Pt
        ImageFile.LOAD_TRUNCATED_IMAGES = True

        try:
            self.logger.info("Started seeding data into cae quality slide")
            self.logger.info("")
            starttime = datetime.now()
            #checking for 'cae quality' window
            if not self.general_input.cae_quality_window_name in ["null","none",""]:
                cae_quality_window_name = self.general_input.cae_quality_window_name
                cae_quality_window_obj = windows.WindowByName(cae_quality_window_name)
                if cae_quality_window_obj:
                    #maximizing the cae quality window
                    utils.MetaCommand('window maximize "{}"'.format(cae_quality_window_name))
                    #iterating through the shapes of the cae quality slide
                    for shape in self.shapes:
                        #image insertion for the shape named "Image 2"
                        if shape.name == "Image 2":
                            #getting "System Energy" plot object to activate
                            plot_id = 0
                            page_id = 0
                            plot = plot2d.Plot(plot_id, cae_quality_window_name, page_id)
                            title = plot.get_title()
                            plot.activate()
                            #changing the layout of the window and showing all the curves
                            utils.MetaCommand('xyplot rlayout "{}" 1'.format(cae_quality_window_name))
                            utils.MetaCommand('xyplot curve select "{}" all'.format(cae_quality_window_name))
                            utils.MetaCommand('xyplot curve visible and "{}" sel'.format(cae_quality_window_name))
                            utils.MetaCommand('xyplot plotoptions value off "{}" {}'.format(cae_quality_window_name, plot_id))
                            #hiding "System damping energy" curve
                            system_damping_energy_curves = plot2d.CurvesByName(cae_quality_window_name, "System damping energy", 1)
                            if system_damping_energy_curves:
                                system_damping_energy_curves[0].hide()
                            #capturing "System Energy" plot image
                            if self.twod_images_report_folder is not None:
                                image_path = os.path.join(self.twod_images_report_folder,cae_quality_window_name+"_"+title.get_text()+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : All except 'System damping energy' curve | SOURCE PLOT : {} | SOURCE WINDOW : {}".format(title.get_text().lower(),cae_quality_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                                #showing legend for "System Energy" plot and moving to left side of the window
                                utils.MetaCommand('xyplot plotoptions legend on "{}" 0'.format(cae_quality_window_name))
                                utils.MetaCommand('xyplot legend hook left "{}" 0'.format(cae_quality_window_name))
                                utils.MetaCommand('xyplot legend hook hout "{}" 0'.format(cae_quality_window_name))
                                utils.MetaCommand('xyplot legend ymove "{}" 0 1.060000'.format(cae_quality_window_name))
                                #capturing "System Energy" plot image along with legend
                                image2_path = os.path.join(self.twod_images_report_folder,cae_quality_window_name+"_"+title.get_text()+"_LEGEND"+".png")
                                utils.MetaCommand('write png "{}"'.format(image2_path))
                                #creating Image object for the above captured image
                                img_2 = Image.open(image2_path)
                                #getting legend attributes of "System Energy" plot
                                legend = plot2d.Legend(plot_id, cae_quality_window_name, page_id)
                                left,top = legend.get_position()
                                width = legend.get_width()
                                height = legend.get_height()
                                #cropping the image to get legend alone
                                img_2 = img_2.crop((left,top,width+8,height+8))
                                img_2.save(image2_path,"PNG")
                                self.logger.info(image2_path)
                                #getting the shape object with name as "Image 1"
                                shape2 = [shape for shape in self.slide.shapes if shape.name == "Image 1"][0]
                                #adding picture based on the above shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.slide.shapes.add_picture(image2_path,shape2.left,shape2.top,width = shape2.width,height = shape2.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                            #reverting the plot layout of the window and deactiving the "System Energy" plot
                            utils.MetaCommand('xyplot rlayout "{}" 2'.format(cae_quality_window_name))
                            plot.deactivate()
                        #image insertion for the shape named "Image 3"
                        elif shape.name == "Image 3":
                            #getting "Added Mass" plot object to activate
                            plot_id = 1
                            page_id=0
                            plot = plot2d.Plot(plot_id, cae_quality_window_name, page_id)
                            title = plot.get_title()
                            plot.activate()
                            #changing the layout of the window,showing all the curves and color update for the visible curves
                            utils.MetaCommand('xyplot rlayout "{}" 1'.format(cae_quality_window_name))
                            utils.MetaCommand('xyplot curve select "{}" all'.format(cae_quality_window_name))
                            utils.MetaCommand('xyplot curve visible and "{}" sel'.format(cae_quality_window_name))
                            utils.MetaCommand('xyplot curve set color "{}" vis LightGreen'.format(cae_quality_window_name))
                            utils.MetaCommand('xyplot plotoptions value off "{}" {}'.format(cae_quality_window_name, plot_id))
                            #capturing "Added Mass" plot image
                            if self.twod_images_report_folder is not None:
                                image_path = os.path.join(self.twod_images_report_folder,cae_quality_window_name+"_"+title.get_text()+".png").replace(" ","_")
                                capture_image_and_resize(image_path,shape.width,shape.height)
                                self.logger.info("--- 2D CURVE IMAGE GENERATOR")
                                self.logger.info("")
                                self.logger.info("CURVES : All | SOURCE PLOT : {} | SOURCE WINDOW : {}".format(title.get_text().lower(),cae_quality_window_name))
                                self.logger.info("OUTPUT IMAGE SIZE (PIXELS) : {}x{}".format(round(shape.width/9525),round(shape.height/9525)))
                                self.logger.info("OUTPUT CURVE IMAGES : ")
                                self.logger.info(image_path)
                                self.logger.info("")
                                #adding picture based on the shape width and height, which will hide the original shape and add a picture shape on top of that
                                picture = self.slide.shapes.add_picture(image_path,shape.left,shape.top,width = shape.width,height = shape.height)
                                picture.crop_left = 0
                                picture.crop_right = 0
                            else:
                                self.logger.info("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.report_directory_key))
                                self.logger.info("")
                            #reverting the plot layout of the window and deactiving the "System Energy" plot
                            utils.MetaCommand('xyplot rlayout "{}" 2'.format(cae_quality_window_name))
                            plot.deactivate()
                        #table population for the shape named "Table 1"
                        elif shape.name == "Table 1":
                            #getting "System Energy" plot object
                            plot_id = 0
                            page_id=0
                            plot = plot2d.Plot(plot_id, cae_quality_window_name, page_id)
                            #getting all the visible curves but skipping System damping energy
                            curvelist = plot.get_curves('all')
                            curvelist = [curve for curve in curvelist if curve.name != "System damping energy"]
                            #getting table object
                            table_obj = shape.table
                            #iterating through all the  curves to populate table with data
                            for index,curve in enumerate(curvelist):
                                #getting min and max values of y values
                                min_y = curve.get_limit_value_y(specifier = 'min')
                                max_y = curve.get_limit_value_y(specifier = 'max')
                                #adding a new row to the table to insert curve data
                                add_row(table_obj)
                                #getting row object and inserting curve name in cell 0
                                row = table_obj.rows[index+1]
                                text_frame_1 = row.cells[0].text_frame
                                font_1 = text_frame_1.paragraphs[0].font
                                font_1.size = Pt(12)
                                text_frame_1.paragraphs[0].text = str(curve.name).replace(" energy","")
                                #inserting max y value of the curve in cell 1
                                text_frame_2 = row.cells[1].text_frame
                                font_2 = text_frame_2.paragraphs[0].font
                                font_2.size = Pt(12)
                                text_frame_2.paragraphs[0].text = "{:.2e}".format(max_y)
                                #inserting min y value of the curve in cell 2
                                text_frame_3 = row.cells[2].text_frame
                                font_3 = text_frame_3.paragraphs[0].font
                                font_3.size = Pt(12)
                                text_frame_3.paragraphs[0].text = "{:.2e}".format(min_y)
                        #table population for the shape named "Table 2"
                        elif shape.name == "Table 2":
                            #getting table object
                            table_obj = shape.table
                            #creating dictionary with data which is used to insert data into table
                            table_value_dict ={"Termination type":self.general_input.termination_type if self.general_input.termination_type not in ["null","none",""] else self.logger.log("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.termination_type_key)),
                                                "Computation time":self.general_input.computation_time if self.general_input.computation_time not in ["null","none",""] else self.logger.log("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.computation_time_key)),
                                                "Core count":self.general_input.core_count if self.general_input.core_count not in ["null","none",""] else self.logger.log("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.core_count_key)),
                                                "Verification mode":self.general_input.verification_mode if self.general_input.verification_mode not in ["null","none",""] else self.logger.log("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.verification_mode_key)),
                                                "Compute cluster":self.general_input.compute_cluster if self.general_input.compute_cluster not in ["null","none",""] else self.logger.log("WARNING : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.compute_cluster_key))}
                            #iterating through above dictionary to insert data into table
                            for index,(item,value) in enumerate(table_value_dict.items()):
                                #adding a new row to the table
                                add_row(table_obj)
                                #getting row object and inserting curve name in cell 0
                                row = table_obj.rows[index+1]
                                text_frame_1 = row.cells[0].text_frame
                                font_1 = text_frame_1.paragraphs[0].font
                                font_1.size = Pt(12)
                                text_frame_1.paragraphs[0].text = item
                                #removing unnecessary string from core count string value
                                if item == "Core count":
                                    value = value.split("with")[1].rstrip()
                                #inserting value in cell 1
                                text_frame_2 = row.cells[1].text_frame
                                font_2 = text_frame_2.paragraphs[0].font
                                font_2.size = Pt(12)
                                text_frame_2.paragraphs[0].text = value if value else ""
                else:
                    self.logger.info("ERROR : 2D METADB does not contain 'Cae Quality' window. Please update.")
            else:
                self.logger.info("ERROR : META 2D variable '{}' is not available or invalid. Please update.".format(GeneralVarInfo.cae_window_key))
            endtime = datetime.now()
        except Exception as e:
            self.logger.exception("Error while seeding data into cae quality slide:\n{}".format(e))
            self.logger.info("")
            return 1
        self.logger.info("Completed seeding data into cae quality slide")
        self.logger.info("Time Taken : {}".format(endtime - starttime))
        self.logger.info("")

        return 0
